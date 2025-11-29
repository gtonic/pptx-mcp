#!/usr/bin/env python3
"""
Tests for the business diagrams module.

Tests the high-level APIs for creating SWOT analysis, Timeline, and Org Chart diagrams.
"""

import sys
import os
sys.path.insert(0, os.path.dirname(os.path.dirname(os.path.abspath(__file__))))

import unittest
from unittest.mock import patch, MagicMock
from pptx import Presentation

from business_diagrams import (
    BusinessDiagramsEngine, 
    SWOT_COLORS, 
    SWOT_COLORS_RGB,
    TIMELINE_COLORS
)


# Slide layout index for blank slides in default PowerPoint templates
BLANK_SLIDE_LAYOUT_INDEX = 6


class TestSWOTAnalysis(unittest.TestCase):
    """Tests for SWOT analysis creation."""
    
    def setUp(self):
        """Set up test fixtures."""
        self.engine = BusinessDiagramsEngine()
        # Create a mock presentation with a blank slide
        self.pres = Presentation()
        self.pres.slides.add_slide(self.pres.slide_layouts[BLANK_SLIDE_LAYOUT_INDEX])
    
    @patch('business_diagrams.presentation_manager')
    def test_create_swot_basic(self, mock_pm):
        """Test basic SWOT creation with all quadrants."""
        mock_pm.get_presentation.return_value = self.pres
        
        result = self.engine.create_swot_analysis(
            slide_index=0,
            strengths=["Strong brand", "Good team"],
            weaknesses=["High costs"],
            opportunities=["New market"],
            threats=["Competition"]
        )
        
        self.assertNotIn("error", result)
        self.assertEqual(result["diagram_type"], "swot_analysis")
        self.assertEqual(result["item_counts"]["strengths"], 2)
        self.assertEqual(result["item_counts"]["weaknesses"], 1)
        self.assertEqual(result["item_counts"]["opportunities"], 1)
        self.assertEqual(result["item_counts"]["threats"], 1)
    
    @patch('business_diagrams.presentation_manager')
    def test_create_swot_with_title(self, mock_pm):
        """Test SWOT creation with a title."""
        mock_pm.get_presentation.return_value = self.pres
        
        result = self.engine.create_swot_analysis(
            slide_index=0,
            strengths=["Item 1"],
            weaknesses=["Item 2"],
            opportunities=["Item 3"],
            threats=["Item 4"],
            title="Company SWOT Analysis"
        )
        
        self.assertNotIn("error", result)
        self.assertEqual(result["diagram_type"], "swot_analysis")
    
    @patch('business_diagrams.presentation_manager')
    def test_create_swot_empty_quadrants(self, mock_pm):
        """Test SWOT creation with some empty quadrants."""
        mock_pm.get_presentation.return_value = self.pres
        
        result = self.engine.create_swot_analysis(
            slide_index=0,
            strengths=["Only strength"],
            weaknesses=[],
            opportunities=[],
            threats=[]
        )
        
        self.assertNotIn("error", result)
        self.assertEqual(result["item_counts"]["strengths"], 1)
        self.assertEqual(result["item_counts"]["weaknesses"], 0)
    
    @patch('business_diagrams.presentation_manager')
    def test_create_swot_invalid_slide(self, mock_pm):
        """Test SWOT creation with invalid slide index."""
        mock_pm.get_presentation.return_value = self.pres
        
        result = self.engine.create_swot_analysis(
            slide_index=99,
            strengths=["Item"],
            weaknesses=["Item"],
            opportunities=["Item"],
            threats=["Item"]
        )
        
        self.assertIn("error", result)
    
    @patch('business_diagrams.presentation_manager')
    def test_create_swot_without_labels(self, mock_pm):
        """Test SWOT creation without category labels."""
        mock_pm.get_presentation.return_value = self.pres
        
        result = self.engine.create_swot_analysis(
            slide_index=0,
            strengths=["Item 1"],
            weaknesses=["Item 2"],
            opportunities=["Item 3"],
            threats=["Item 4"],
            show_labels=False
        )
        
        self.assertNotIn("error", result)


class TestTimeline(unittest.TestCase):
    """Tests for timeline creation."""
    
    def setUp(self):
        """Set up test fixtures."""
        self.engine = BusinessDiagramsEngine()
        self.pres = Presentation()
        self.pres.slides.add_slide(self.pres.slide_layouts[BLANK_SLIDE_LAYOUT_INDEX])
    
    @patch('business_diagrams.presentation_manager')
    def test_create_horizontal_timeline(self, mock_pm):
        """Test horizontal timeline creation."""
        mock_pm.get_presentation.return_value = self.pres
        
        events = [
            {"label": "Start", "date": "Jan 2024"},
            {"label": "Middle", "date": "Mar 2024"},
            {"label": "End", "date": "Jun 2024"}
        ]
        
        result = self.engine.create_timeline(
            slide_index=0,
            events=events,
            direction="horizontal"
        )
        
        self.assertNotIn("error", result)
        self.assertEqual(result["diagram_type"], "timeline")
        self.assertEqual(result["direction"], "horizontal")
        self.assertEqual(result["event_count"], 3)
    
    @patch('business_diagrams.presentation_manager')
    def test_create_vertical_timeline(self, mock_pm):
        """Test vertical timeline creation."""
        mock_pm.get_presentation.return_value = self.pres
        
        events = [
            {"label": "Phase 1"},
            {"label": "Phase 2"},
            {"label": "Phase 3"}
        ]
        
        result = self.engine.create_timeline(
            slide_index=0,
            events=events,
            direction="vertical"
        )
        
        self.assertNotIn("error", result)
        self.assertEqual(result["direction"], "vertical")
    
    @patch('business_diagrams.presentation_manager')
    def test_create_timeline_with_colors(self, mock_pm):
        """Test timeline with custom event colors."""
        mock_pm.get_presentation.return_value = self.pres
        
        events = [
            {"label": "Start", "color": "success"},
            {"label": "Milestone", "color": [255, 0, 0]},
            {"label": "End", "color": "accent"}
        ]
        
        result = self.engine.create_timeline(
            slide_index=0,
            events=events
        )
        
        self.assertNotIn("error", result)
    
    @patch('business_diagrams.presentation_manager')
    def test_create_timeline_with_descriptions(self, mock_pm):
        """Test timeline with event descriptions."""
        mock_pm.get_presentation.return_value = self.pres
        
        events = [
            {"label": "Kickoff", "date": "Jan", "description": "Project begins"},
            {"label": "Launch", "date": "Dec", "description": "Product goes live"}
        ]
        
        result = self.engine.create_timeline(
            slide_index=0,
            events=events
        )
        
        self.assertNotIn("error", result)
    
    @patch('business_diagrams.presentation_manager')
    def test_create_timeline_no_events(self, mock_pm):
        """Test timeline with no events."""
        mock_pm.get_presentation.return_value = self.pres
        
        result = self.engine.create_timeline(
            slide_index=0,
            events=[]
        )
        
        self.assertIn("error", result)
    
    @patch('business_diagrams.presentation_manager')
    def test_create_timeline_without_connector(self, mock_pm):
        """Test timeline without connector line."""
        mock_pm.get_presentation.return_value = self.pres
        
        events = [{"label": "Event 1"}, {"label": "Event 2"}]
        
        result = self.engine.create_timeline(
            slide_index=0,
            events=events,
            show_connector=False
        )
        
        self.assertNotIn("error", result)


class TestOrgChart(unittest.TestCase):
    """Tests for organization chart creation."""
    
    def setUp(self):
        """Set up test fixtures."""
        self.engine = BusinessDiagramsEngine()
        self.pres = Presentation()
        self.pres.slides.add_slide(self.pres.slide_layouts[BLANK_SLIDE_LAYOUT_INDEX])
    
    @patch('business_diagrams.presentation_manager')
    @patch('business_diagrams.layout_manager')
    def test_create_simple_org_chart(self, mock_lm, mock_pm):
        """Test simple org chart with one level."""
        mock_pm.get_presentation.return_value = self.pres
        mock_lm.create_hierarchy_layout.return_value = {
            "message": "Created hierarchy",
            "levels": 1,
            "shapes": []
        }
        mock_lm._get_slide_bounds.return_value = MagicMock(
            left=0.5, top=1.0, width=9.0, height=5.5
        )
        
        root = {
            "name": "CEO",
            "title": "Chief Executive"
        }
        
        result = self.engine.create_org_chart(
            slide_index=0,
            root=root
        )
        
        self.assertNotIn("error", result)
        self.assertEqual(result["diagram_type"], "org_chart")
    
    @patch('business_diagrams.presentation_manager')
    @patch('business_diagrams.layout_manager')
    def test_create_org_chart_with_children(self, mock_lm, mock_pm):
        """Test org chart with multiple levels."""
        mock_pm.get_presentation.return_value = self.pres
        mock_lm.create_hierarchy_layout.return_value = {
            "message": "Created hierarchy",
            "levels": 3,
            "shapes": []
        }
        mock_lm._get_slide_bounds.return_value = MagicMock(
            left=0.5, top=1.0, width=9.0, height=5.5
        )
        
        root = {
            "name": "CEO",
            "title": "Chief Executive",
            "children": [
                {
                    "name": "VP Sales",
                    "children": [
                        {"name": "Sales Rep 1"},
                        {"name": "Sales Rep 2"}
                    ]
                },
                {
                    "name": "VP Engineering",
                    "children": [
                        {"name": "Dev 1"}
                    ]
                }
            ]
        }
        
        result = self.engine.create_org_chart(
            slide_index=0,
            root=root
        )
        
        self.assertNotIn("error", result)
    
    @patch('business_diagrams.presentation_manager')
    @patch('business_diagrams.layout_manager')
    def test_create_org_chart_compact(self, mock_lm, mock_pm):
        """Test org chart with compact mode."""
        mock_pm.get_presentation.return_value = self.pres
        mock_lm.create_hierarchy_layout.return_value = {
            "message": "Created hierarchy",
            "levels": 2,
            "shapes": []
        }
        mock_lm._get_slide_bounds.return_value = MagicMock(
            left=0.5, top=1.0, width=9.0, height=5.5
        )
        
        root = {
            "name": "Manager",
            "children": [{"name": "Report 1"}, {"name": "Report 2"}]
        }
        
        result = self.engine.create_org_chart(
            slide_index=0,
            root=root,
            compact=True
        )
        
        self.assertNotIn("error", result)
    
    def test_transform_org_node_basic(self):
        """Test org node transformation."""
        node = {
            "name": "John Doe",
            "title": "Manager"
        }
        
        result = self.engine._transform_org_node(node)
        
        self.assertIn("content", result)
        self.assertIn("John Doe", result["content"])
        self.assertIn("Manager", result["content"])
        self.assertEqual(result["element_type"], "shape")
        self.assertEqual(result["shape_type"], "rounded_rectangle")
    
    def test_transform_org_node_with_children(self):
        """Test org node transformation with children."""
        node = {
            "name": "Parent",
            "children": [
                {"name": "Child 1"},
                {"name": "Child 2"}
            ]
        }
        
        result = self.engine._transform_org_node(node)
        
        self.assertIn("children", result)
        self.assertEqual(len(result["children"]), 2)
    
    def test_transform_org_node_with_color(self):
        """Test org node transformation with custom color."""
        node = {
            "name": "Person",
            "color": [100, 150, 200]
        }
        
        result = self.engine._transform_org_node(node)
        
        self.assertEqual(result["fill_color"], [100, 150, 200])


class TestColorResolution(unittest.TestCase):
    """Tests for color resolution in business diagrams."""
    
    def test_swot_colors_defined(self):
        """Test that SWOT colors are properly defined."""
        self.assertIn("strengths", SWOT_COLORS)
        self.assertIn("weaknesses", SWOT_COLORS)
        self.assertIn("opportunities", SWOT_COLORS)
        self.assertIn("threats", SWOT_COLORS)
    
    def test_swot_colors_rgb_fallback(self):
        """Test that RGB fallback colors are defined."""
        for key in SWOT_COLORS:
            self.assertIn(key, SWOT_COLORS_RGB)
            rgb = SWOT_COLORS_RGB[key]
            self.assertEqual(len(rgb), 3)
            for val in rgb:
                self.assertGreaterEqual(val, 0)
                self.assertLessEqual(val, 255)
    
    def test_timeline_colors_defined(self):
        """Test that timeline colors are properly defined."""
        self.assertIn("connector", TIMELINE_COLORS)
        self.assertIn("event_default", TIMELINE_COLORS)


if __name__ == '__main__':
    unittest.main()
