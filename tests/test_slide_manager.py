#!/usr/bin/env python3
"""
Tests for the slide manager module.

Tests slide operations including adding slides, textboxes, shapes, charts, tables,
images, bullet points, and auto-fit text functionality.
"""

import sys
import os
sys.path.insert(0, os.path.dirname(os.path.dirname(os.path.abspath(__file__))))

import unittest
from unittest.mock import patch, MagicMock
from pptx import Presentation

from slide_manager import SlideManager, slide_manager


# Slide layout index for blank slides in default PowerPoint templates
BLANK_SLIDE_LAYOUT_INDEX = 6


class TestResolveColorWithDefault(unittest.TestCase):
    """Tests for the _resolve_color_with_default helper function."""
    
    def setUp(self):
        """Set up test fixtures."""
        self.manager = SlideManager()
    
    @patch('slide_manager.template_manager')
    def test_returns_rgb_list_when_provided(self, mock_tm):
        """Test that RGB list is returned when provided."""
        mock_tm.resolve_color.return_value = [255, 0, 0]
        
        result = self.manager._resolve_color_with_default([255, 0, 0])
        
        self.assertEqual(result, [255, 0, 0])
    
    @patch('slide_manager.template_manager')
    def test_returns_resolved_semantic_color(self, mock_tm):
        """Test that semantic color tags are resolved."""
        mock_tm.resolve_color.return_value = [0, 128, 255]
        
        result = self.manager._resolve_color_with_default("accent")
        
        mock_tm.resolve_color.assert_called_once_with("accent")
        self.assertEqual(result, [0, 128, 255])
    
    @patch('slide_manager.template_manager')
    def test_uses_accent_1_default_when_color_is_none(self, mock_tm):
        """Test that accent_1 is used as default when no color is provided."""
        mock_tm.resolve_color.return_value = None
        mock_tm.get_default_color_settings.return_value = {
            "accent_1": (79, 129, 189),
            "text_1": (0, 0, 0)
        }
        
        result = self.manager._resolve_color_with_default(None)
        
        self.assertEqual(result, [79, 129, 189])
    
    @patch('slide_manager.template_manager')
    def test_uses_custom_default_key(self, mock_tm):
        """Test that custom default key is used when specified."""
        mock_tm.resolve_color.return_value = None
        mock_tm.get_default_color_settings.return_value = {
            "accent_1": (79, 129, 189),
            "text_1": (0, 0, 0)
        }
        
        result = self.manager._resolve_color_with_default(None, default_key="text_1")
        
        self.assertEqual(result, [0, 0, 0])
    
    @patch('slide_manager.template_manager')
    def test_returns_none_when_no_default_available(self, mock_tm):
        """Test that None is returned when no default is available."""
        mock_tm.resolve_color.return_value = None
        mock_tm.get_default_color_settings.return_value = {}
        
        result = self.manager._resolve_color_with_default(None)
        
        self.assertIsNone(result)
    
    @patch('slide_manager.template_manager')
    def test_returns_none_for_missing_default_key(self, mock_tm):
        """Test that None is returned when default key doesn't exist."""
        mock_tm.resolve_color.return_value = None
        mock_tm.get_default_color_settings.return_value = {
            "accent_1": (79, 129, 189)
        }
        
        result = self.manager._resolve_color_with_default(None, default_key="nonexistent_key")
        
        self.assertIsNone(result)
    
    @patch('slide_manager.template_manager')
    def test_converts_tuple_to_list(self, mock_tm):
        """Test that tuple color values are converted to lists."""
        mock_tm.resolve_color.return_value = None
        mock_tm.get_default_color_settings.return_value = {
            "accent_1": (100, 150, 200)  # Tuple, as returned by get_default_color_settings
        }
        
        result = self.manager._resolve_color_with_default(None)
        
        self.assertIsInstance(result, list)
        self.assertEqual(result, [100, 150, 200])
    
    @patch('slide_manager.template_manager')
    def test_does_not_use_default_when_color_resolved(self, mock_tm):
        """Test that default is not used when color is successfully resolved."""
        mock_tm.resolve_color.return_value = [255, 128, 0]
        mock_tm.get_default_color_settings.return_value = {
            "accent_1": (79, 129, 189)
        }
        
        result = self.manager._resolve_color_with_default("warning")
        
        mock_tm.get_default_color_settings.assert_not_called()
        self.assertEqual(result, [255, 128, 0])


class TestResolveColor(unittest.TestCase):
    """Tests for the _resolve_color method."""
    
    def setUp(self):
        """Set up test fixtures."""
        self.manager = SlideManager()
    
    @patch('slide_manager.template_manager')
    def test_delegates_to_template_manager(self, mock_tm):
        """Test that _resolve_color delegates to template_manager."""
        mock_tm.resolve_color.return_value = [128, 64, 32]
        
        result = self.manager._resolve_color("test_color")
        
        mock_tm.resolve_color.assert_called_once_with("test_color")
        self.assertEqual(result, [128, 64, 32])
    
    @patch('slide_manager.template_manager')
    def test_returns_none_for_none_input(self, mock_tm):
        """Test that None input returns None."""
        mock_tm.resolve_color.return_value = None
        
        result = self.manager._resolve_color(None)
        
        self.assertIsNone(result)


class TestSlideManagerInitialization(unittest.TestCase):
    """Tests for SlideManager initialization."""
    
    def test_global_instance_exists(self):
        """Test that global slide_manager instance exists."""
        self.assertIsInstance(slide_manager, SlideManager)


class TestAddSlide(unittest.TestCase):
    """Tests for add_slide method."""
    
    def setUp(self):
        """Set up test fixtures."""
        self.manager = SlideManager()
        self.pres = Presentation()
    
    @patch('slide_manager.presentation_manager')
    def test_add_slide_basic(self, mock_pm):
        """Test basic slide addition."""
        mock_pm.get_presentation.return_value = self.pres
        
        result = self.manager.add_slide(layout_index=BLANK_SLIDE_LAYOUT_INDEX)
        
        self.assertNotIn("error", result)
        self.assertIn("slide_index", result)
        self.assertEqual(result["slide_index"], 0)
    
    @patch('slide_manager.presentation_manager')
    def test_add_slide_with_title(self, mock_pm):
        """Test adding slide with title."""
        mock_pm.get_presentation.return_value = self.pres
        
        result = self.manager.add_slide(
            layout_index=0,  # Title slide layout
            title="Test Title"
        )
        
        self.assertNotIn("error", result)
    
    @patch('slide_manager.presentation_manager')
    def test_add_slide_invalid_layout(self, mock_pm):
        """Test adding slide with invalid layout index."""
        mock_pm.get_presentation.return_value = self.pres
        
        result = self.manager.add_slide(layout_index=999)
        
        self.assertIn("error", result)
        self.assertIn("available_layouts", result)
    
    @patch('slide_manager.presentation_manager')
    def test_add_slide_negative_layout(self, mock_pm):
        """Test adding slide with negative layout index."""
        mock_pm.get_presentation.return_value = self.pres
        
        result = self.manager.add_slide(layout_index=-1)
        
        self.assertIn("error", result)
    
    @patch('slide_manager.presentation_manager')
    def test_add_multiple_slides(self, mock_pm):
        """Test adding multiple slides."""
        mock_pm.get_presentation.return_value = self.pres
        
        result1 = self.manager.add_slide(layout_index=BLANK_SLIDE_LAYOUT_INDEX)
        result2 = self.manager.add_slide(layout_index=BLANK_SLIDE_LAYOUT_INDEX)
        
        self.assertEqual(result1["slide_index"], 0)
        self.assertEqual(result2["slide_index"], 1)


class TestAddTextbox(unittest.TestCase):
    """Tests for add_textbox method."""
    
    def setUp(self):
        """Set up test fixtures."""
        self.manager = SlideManager()
        self.pres = Presentation()
        self.pres.slides.add_slide(self.pres.slide_layouts[BLANK_SLIDE_LAYOUT_INDEX])
    
    @patch('slide_manager.template_manager')
    @patch('slide_manager.presentation_manager')
    def test_add_textbox_basic(self, mock_pm, mock_tm):
        """Test basic textbox addition."""
        mock_pm.get_presentation.return_value = self.pres
        mock_tm.resolve_color.return_value = None
        mock_tm.get_default_color_settings.return_value = {"accent_1": (79, 129, 189)}
        mock_tm.resolve_font.return_value = {}
        
        result = self.manager.add_textbox(
            slide_index=0,
            left=1.0,
            top=1.0,
            width=2.0,
            height=1.0,
            text="Test Text"
        )
        
        self.assertNotIn("error", result)
        self.assertIn("shape_index", result)
    
    @patch('slide_manager.template_manager')
    @patch('slide_manager.presentation_manager')
    def test_add_textbox_with_formatting(self, mock_pm, mock_tm):
        """Test textbox addition with formatting options."""
        mock_pm.get_presentation.return_value = self.pres
        mock_tm.resolve_color.return_value = [255, 0, 0]
        mock_tm.resolve_font.return_value = {"font_name": "Arial", "font_size": 18}
        
        result = self.manager.add_textbox(
            slide_index=0,
            left=1.0,
            top=1.0,
            width=2.0,
            height=1.0,
            text="Formatted Text",
            font_size=18,
            font_name="Arial",
            bold=True,
            color=[255, 0, 0],
            alignment="center"
        )
        
        self.assertNotIn("error", result)
    
    @patch('slide_manager.template_manager')
    @patch('slide_manager.presentation_manager')
    def test_add_textbox_invalid_slide(self, mock_pm, mock_tm):
        """Test textbox addition with invalid slide index."""
        mock_pm.get_presentation.return_value = self.pres
        
        result = self.manager.add_textbox(
            slide_index=99,
            left=1.0,
            top=1.0,
            width=2.0,
            height=1.0,
            text="Test"
        )
        
        self.assertIn("error", result)
    
    @patch('slide_manager.template_manager')
    @patch('slide_manager.presentation_manager')
    def test_add_textbox_semantic_color(self, mock_pm, mock_tm):
        """Test textbox with semantic color tag."""
        mock_pm.get_presentation.return_value = self.pres
        mock_tm.resolve_color.return_value = [0, 128, 255]
        mock_tm.resolve_font.return_value = {}
        
        result = self.manager.add_textbox(
            slide_index=0,
            left=1.0,
            top=1.0,
            width=2.0,
            height=1.0,
            text="Semantic Color",
            color="accent"
        )
        
        self.assertNotIn("error", result)
        mock_tm.resolve_color.assert_called_with("accent")


class TestAddShape(unittest.TestCase):
    """Tests for add_shape method."""
    
    def setUp(self):
        """Set up test fixtures."""
        self.manager = SlideManager()
        self.pres = Presentation()
        self.pres.slides.add_slide(self.pres.slide_layouts[BLANK_SLIDE_LAYOUT_INDEX])
    
    @patch('slide_manager.template_manager')
    @patch('slide_manager.presentation_manager')
    def test_add_shape_rectangle(self, mock_pm, mock_tm):
        """Test adding rectangle shape."""
        mock_pm.get_presentation.return_value = self.pres
        mock_tm.resolve_color.return_value = None
        mock_tm.get_default_color_settings.return_value = {
            "accent_1": (79, 129, 189),
            "text_1": (0, 0, 0)
        }
        
        result = self.manager.add_shape(
            slide_index=0,
            shape_type="rectangle",
            left=1.0,
            top=1.0,
            width=2.0,
            height=1.0
        )
        
        self.assertNotIn("error", result)
        self.assertIn("shape_index", result)
    
    @patch('slide_manager.template_manager')
    @patch('slide_manager.presentation_manager')
    def test_add_shape_with_colors(self, mock_pm, mock_tm):
        """Test adding shape with fill and line colors."""
        mock_pm.get_presentation.return_value = self.pres
        mock_tm.resolve_color.return_value = [255, 0, 0]
        mock_tm.get_default_color_settings.return_value = {
            "accent_1": (79, 129, 189),
            "text_1": (0, 0, 0)
        }
        
        result = self.manager.add_shape(
            slide_index=0,
            shape_type="oval",
            left=1.0,
            top=1.0,
            width=2.0,
            height=2.0,
            fill_color=[255, 0, 0],
            line_color=[0, 0, 0],
            line_width=2.0
        )
        
        self.assertNotIn("error", result)
    
    @patch('slide_manager.template_manager')
    @patch('slide_manager.presentation_manager')
    def test_add_shape_invalid_slide(self, mock_pm, mock_tm):
        """Test adding shape with invalid slide index."""
        mock_pm.get_presentation.return_value = self.pres
        
        result = self.manager.add_shape(
            slide_index=99,
            shape_type="rectangle",
            left=1.0,
            top=1.0,
            width=2.0,
            height=1.0
        )
        
        self.assertIn("error", result)


class TestAddLine(unittest.TestCase):
    """Tests for add_line method."""
    
    def setUp(self):
        """Set up test fixtures."""
        self.manager = SlideManager()
        self.pres = Presentation()
        self.pres.slides.add_slide(self.pres.slide_layouts[BLANK_SLIDE_LAYOUT_INDEX])
    
    @patch('slide_manager.template_manager')
    @patch('slide_manager.presentation_manager')
    def test_add_line_basic(self, mock_pm, mock_tm):
        """Test adding basic line."""
        mock_pm.get_presentation.return_value = self.pres
        mock_tm.resolve_color.return_value = None
        mock_tm.get_default_color_settings.return_value = {"text_1": (0, 0, 0)}
        
        result = self.manager.add_line(
            slide_index=0,
            x1=1.0,
            y1=1.0,
            x2=3.0,
            y2=1.0
        )
        
        self.assertNotIn("error", result)
        self.assertIn("shape_index", result)
    
    @patch('slide_manager.template_manager')
    @patch('slide_manager.presentation_manager')
    def test_add_line_with_style(self, mock_pm, mock_tm):
        """Test adding line with custom style."""
        mock_pm.get_presentation.return_value = self.pres
        mock_tm.resolve_color.return_value = [255, 0, 0]
        mock_tm.get_default_color_settings.return_value = {"text_1": (0, 0, 0)}
        
        result = self.manager.add_line(
            slide_index=0,
            x1=1.0,
            y1=1.0,
            x2=3.0,
            y2=3.0,
            line_color=[255, 0, 0],
            line_width=3.0
        )
        
        self.assertNotIn("error", result)
    
    @patch('slide_manager.template_manager')
    @patch('slide_manager.presentation_manager')
    def test_add_line_invalid_slide(self, mock_pm, mock_tm):
        """Test adding line with invalid slide index."""
        mock_pm.get_presentation.return_value = self.pres
        
        result = self.manager.add_line(
            slide_index=99,
            x1=1.0,
            y1=1.0,
            x2=3.0,
            y2=1.0
        )
        
        self.assertIn("error", result)


class TestAddChart(unittest.TestCase):
    """Tests for add_chart method."""
    
    def setUp(self):
        """Set up test fixtures."""
        self.manager = SlideManager()
        self.pres = Presentation()
        self.pres.slides.add_slide(self.pres.slide_layouts[BLANK_SLIDE_LAYOUT_INDEX])
    
    @patch('slide_manager.presentation_manager')
    def test_add_chart_column(self, mock_pm):
        """Test adding column chart."""
        mock_pm.get_presentation.return_value = self.pres
        
        chart_data = {
            "categories": ["Q1", "Q2", "Q3", "Q4"],
            "series": [
                {"name": "Sales", "values": [100, 150, 120, 180]}
            ]
        }
        
        result = self.manager.add_chart(
            slide_index=0,
            chart_type="column",
            left=1.0,
            top=1.0,
            width=6.0,
            height=4.0,
            data=chart_data
        )
        
        self.assertNotIn("error", result)
        self.assertIn("shape_index", result)
    
    @patch('slide_manager.presentation_manager')
    def test_add_chart_pie(self, mock_pm):
        """Test adding pie chart."""
        mock_pm.get_presentation.return_value = self.pres
        
        chart_data = {
            "categories": ["A", "B", "C"],
            "series": [
                {"name": "Distribution", "values": [30, 50, 20]}
            ]
        }
        
        result = self.manager.add_chart(
            slide_index=0,
            chart_type="pie",
            left=1.0,
            top=1.0,
            width=5.0,
            height=4.0,
            data=chart_data
        )
        
        self.assertNotIn("error", result)
    
    @patch('slide_manager.presentation_manager')
    def test_add_chart_line(self, mock_pm):
        """Test adding line chart."""
        mock_pm.get_presentation.return_value = self.pres
        
        chart_data = {
            "categories": ["Jan", "Feb", "Mar"],
            "series": [
                {"name": "Revenue", "values": [100, 120, 140]},
                {"name": "Expenses", "values": [80, 90, 100]}
            ]
        }
        
        result = self.manager.add_chart(
            slide_index=0,
            chart_type="line",
            left=1.0,
            top=1.0,
            width=6.0,
            height=4.0,
            data=chart_data
        )
        
        self.assertNotIn("error", result)
    
    @patch('slide_manager.presentation_manager')
    def test_add_chart_invalid_slide(self, mock_pm):
        """Test adding chart with invalid slide index."""
        mock_pm.get_presentation.return_value = self.pres
        
        result = self.manager.add_chart(
            slide_index=99,
            chart_type="column",
            left=1.0,
            top=1.0,
            width=6.0,
            height=4.0,
            data={"categories": [], "series": []}
        )
        
        self.assertIn("error", result)


class TestAddTable(unittest.TestCase):
    """Tests for add_table method."""
    
    def setUp(self):
        """Set up test fixtures."""
        self.manager = SlideManager()
        self.pres = Presentation()
        self.pres.slides.add_slide(self.pres.slide_layouts[BLANK_SLIDE_LAYOUT_INDEX])
    
    @patch('slide_manager.presentation_manager')
    def test_add_table_basic(self, mock_pm):
        """Test adding basic table."""
        mock_pm.get_presentation.return_value = self.pres
        
        result = self.manager.add_table(
            slide_index=0,
            left=1.0,
            top=1.0,
            rows=3,
            cols=3
        )
        
        self.assertNotIn("error", result)
        self.assertIn("shape_index", result)
    
    @patch('slide_manager.presentation_manager')
    def test_add_table_with_data(self, mock_pm):
        """Test adding table with data."""
        mock_pm.get_presentation.return_value = self.pres
        
        table_data = [
            ["Header 1", "Header 2", "Header 3"],
            ["Row 1 Col 1", "Row 1 Col 2", "Row 1 Col 3"],
            ["Row 2 Col 1", "Row 2 Col 2", "Row 2 Col 3"]
        ]
        
        result = self.manager.add_table(
            slide_index=0,
            left=1.0,
            top=1.0,
            rows=3,
            cols=3,
            data=table_data
        )
        
        self.assertNotIn("error", result)
    
    @patch('slide_manager.presentation_manager')
    def test_add_table_invalid_slide(self, mock_pm):
        """Test adding table with invalid slide index."""
        mock_pm.get_presentation.return_value = self.pres
        
        result = self.manager.add_table(
            slide_index=99,
            left=1.0,
            top=1.0,
            rows=2,
            cols=2
        )
        
        self.assertIn("error", result)


class TestAddImage(unittest.TestCase):
    """Tests for add_image method."""
    
    def setUp(self):
        """Set up test fixtures."""
        self.manager = SlideManager()
        self.pres = Presentation()
        self.pres.slides.add_slide(self.pres.slide_layouts[BLANK_SLIDE_LAYOUT_INDEX])
    
    @patch('slide_manager.presentation_manager')
    def test_add_image_file_not_found(self, mock_pm):
        """Test adding image that doesn't exist."""
        mock_pm.get_presentation.return_value = self.pres
        
        result = self.manager.add_image(
            slide_index=0,
            image_path="/data/nonexistent.png",
            left=1.0,
            top=1.0
        )
        
        self.assertIn("error", result)
    
    @patch('slide_manager.presentation_manager')
    def test_add_image_invalid_slide(self, mock_pm):
        """Test adding image with invalid slide index."""
        mock_pm.get_presentation.return_value = self.pres
        
        result = self.manager.add_image(
            slide_index=99,
            image_path="/data/test.png",
            left=1.0,
            top=1.0
        )
        
        self.assertIn("error", result)


class TestAddBulletPoints(unittest.TestCase):
    """Tests for add_bullet_points method."""
    
    def setUp(self):
        """Set up test fixtures."""
        self.manager = SlideManager()
        self.pres = Presentation()
        # Use a layout with content placeholder
        self.pres.slides.add_slide(self.pres.slide_layouts[1])  # Title and Content layout
    
    @patch('slide_manager.presentation_manager')
    def test_add_bullet_points_basic(self, mock_pm):
        """Test adding basic bullet points."""
        mock_pm.get_presentation.return_value = self.pres
        
        result = self.manager.add_bullet_points(
            slide_index=0,
            placeholder_idx=1,
            bullet_points=["Point 1", "Point 2", "Point 3"]
        )
        
        self.assertNotIn("error", result)
        self.assertEqual(result["placeholder_index"], 1)
    
    @patch('slide_manager.presentation_manager')
    def test_add_bullet_points_with_font_size(self, mock_pm):
        """Test adding bullet points with custom font size."""
        mock_pm.get_presentation.return_value = self.pres
        
        result = self.manager.add_bullet_points(
            slide_index=0,
            placeholder_idx=1,
            bullet_points=["Large text"],
            font_size=24
        )
        
        self.assertNotIn("error", result)
    
    @patch('slide_manager.presentation_manager')
    def test_add_bullet_points_invalid_slide(self, mock_pm):
        """Test adding bullet points with invalid slide index."""
        mock_pm.get_presentation.return_value = self.pres
        
        result = self.manager.add_bullet_points(
            slide_index=99,
            placeholder_idx=1,
            bullet_points=["Test"]
        )
        
        self.assertIn("error", result)


class TestAddAutoFitText(unittest.TestCase):
    """Tests for add_auto_fit_text method."""
    
    def setUp(self):
        """Set up test fixtures."""
        self.manager = SlideManager()
        self.pres = Presentation()
        self.pres.slides.add_slide(self.pres.slide_layouts[BLANK_SLIDE_LAYOUT_INDEX])
    
    @patch('slide_manager.template_manager')
    @patch('slide_manager.presentation_manager')
    def test_add_auto_fit_text_basic(self, mock_pm, mock_tm):
        """Test adding auto-fit text with default settings."""
        mock_pm.get_presentation.return_value = self.pres
        mock_tm.resolve_color.return_value = None
        mock_tm.get_default_color_settings.return_value = {"accent_1": (79, 129, 189)}
        mock_tm.resolve_font.return_value = {}
        
        result = self.manager.add_auto_fit_text(
            slide_index=0,
            left=1.0,
            top=1.0,
            width=4.0,
            height=3.0,
            text="Short text that fits easily."
        )
        
        self.assertNotIn("error", result)
        self.assertIn("strategy_used", result)
    
    @patch('slide_manager.template_manager')
    @patch('slide_manager.presentation_manager')
    def test_add_auto_fit_text_with_strategy(self, mock_pm, mock_tm):
        """Test adding auto-fit text with specific strategy."""
        mock_pm.get_presentation.return_value = self.pres
        mock_tm.resolve_color.return_value = None
        mock_tm.get_default_color_settings.return_value = {"accent_1": (79, 129, 189)}
        mock_tm.resolve_font.return_value = {}
        
        result = self.manager.add_auto_fit_text(
            slide_index=0,
            left=1.0,
            top=1.0,
            width=4.0,
            height=3.0,
            text="Some longer text content.",
            strategy="shrink_font"
        )
        
        self.assertNotIn("error", result)
    
    @patch('slide_manager.template_manager')
    @patch('slide_manager.presentation_manager')
    def test_add_auto_fit_text_long_content(self, mock_pm, mock_tm):
        """Test adding auto-fit text with long content."""
        mock_pm.get_presentation.return_value = self.pres
        mock_tm.resolve_color.return_value = None
        mock_tm.get_default_color_settings.return_value = {"accent_1": (79, 129, 189)}
        mock_tm.resolve_font.return_value = {}
        
        long_text = "This is a very long text. " * 50
        
        result = self.manager.add_auto_fit_text(
            slide_index=0,
            left=1.0,
            top=1.0,
            width=4.0,
            height=2.0,
            text=long_text,
            strategy="smart"
        )
        
        self.assertNotIn("error", result)
    
    @patch('slide_manager.template_manager')
    @patch('slide_manager.presentation_manager')
    def test_add_auto_fit_text_invalid_slide(self, mock_pm, mock_tm):
        """Test adding auto-fit text with invalid slide index."""
        mock_pm.get_presentation.return_value = self.pres
        
        result = self.manager.add_auto_fit_text(
            slide_index=99,
            left=1.0,
            top=1.0,
            width=4.0,
            height=3.0,
            text="Test"
        )
        
        self.assertIn("error", result)


if __name__ == '__main__':
    unittest.main()
