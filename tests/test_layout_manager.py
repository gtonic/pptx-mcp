#!/usr/bin/env python3
"""
Tests for the layout manager module.

Tests the high-level layout engine for automatic element positioning,
including grid, list, hierarchy, and flow layouts.
"""

import sys
import os
sys.path.insert(0, os.path.dirname(os.path.dirname(os.path.abspath(__file__))))

import unittest
from unittest.mock import patch, MagicMock
from pptx import Presentation

from layout_manager import (
    LayoutEngine, LayoutBounds, LayoutElement,
    LayoutType, Alignment, Direction,
    layout_manager
)


# Slide layout index for blank slides in default PowerPoint templates
BLANK_SLIDE_LAYOUT_INDEX = 6


class TestLayoutEngineInitialization(unittest.TestCase):
    """Tests for LayoutEngine initialization."""
    
    def test_initialization(self):
        """Test that LayoutEngine initializes correctly."""
        engine = LayoutEngine()
        self.assertIsNotNone(engine.default_bounds)
        self.assertIsInstance(engine.default_bounds, LayoutBounds)
    
    def test_global_instance_exists(self):
        """Test that global layout_manager instance exists."""
        self.assertIsInstance(layout_manager, LayoutEngine)


class TestLayoutBounds(unittest.TestCase):
    """Tests for LayoutBounds dataclass."""
    
    def test_default_values(self):
        """Test default values for LayoutBounds."""
        bounds = LayoutBounds()
        self.assertEqual(bounds.left, 0.5)
        self.assertEqual(bounds.top, 1.0)
        self.assertEqual(bounds.width, 9.0)
        self.assertEqual(bounds.height, 5.5)
    
    def test_from_slide(self):
        """Test creating bounds from slide dimensions."""
        bounds = LayoutBounds.from_slide(
            slide_width=10.0,
            slide_height=7.5,
            margin=0.5,
            title_height=1.0
        )
        
        self.assertEqual(bounds.left, 0.5)
        self.assertEqual(bounds.top, 1.5)  # margin + title_height
        self.assertEqual(bounds.width, 9.0)  # slide_width - 2*margin
        self.assertEqual(bounds.height, 5.5)  # slide_height - 2*margin - title_height
    
    def test_from_slide_custom_dimensions(self):
        """Test creating bounds with custom slide dimensions."""
        bounds = LayoutBounds.from_slide(
            slide_width=13.333,  # Widescreen
            slide_height=7.5,
            margin=0.75,
            title_height=1.25
        )
        
        self.assertEqual(bounds.left, 0.75)
        self.assertEqual(bounds.top, 2.0)
        self.assertAlmostEqual(bounds.width, 11.833, places=3)


class TestLayoutElement(unittest.TestCase):
    """Tests for LayoutElement dataclass."""
    
    def test_default_values(self):
        """Test default values for LayoutElement."""
        elem = LayoutElement(content="Test")
        
        self.assertEqual(elem.content, "Test")
        self.assertEqual(elem.element_type, "textbox")
        self.assertIsNone(elem.shape_type)
        self.assertEqual(elem.alignment, "center")
        self.assertEqual(elem.left, 0.0)
        self.assertEqual(elem.top, 0.0)
        self.assertEqual(elem.width, 1.0)
        self.assertEqual(elem.height, 1.0)
    
    def test_custom_values(self):
        """Test creating LayoutElement with custom values."""
        elem = LayoutElement(
            content="Custom",
            element_type="shape",
            shape_type="rectangle",
            fill_color=[255, 0, 0],
            font_size=18,
            bold=True
        )
        
        self.assertEqual(elem.content, "Custom")
        self.assertEqual(elem.element_type, "shape")
        self.assertEqual(elem.shape_type, "rectangle")
        self.assertEqual(elem.fill_color, [255, 0, 0])
        self.assertEqual(elem.font_size, 18)
        self.assertTrue(elem.bold)


class TestCreateGridLayout(unittest.TestCase):
    """Tests for create_grid_layout method."""
    
    def setUp(self):
        """Set up test fixtures."""
        self.engine = LayoutEngine()
        self.pres = Presentation()
        self.pres.slides.add_slide(self.pres.slide_layouts[BLANK_SLIDE_LAYOUT_INDEX])
    
    @patch('layout_manager.presentation_manager')
    def test_create_grid_layout_basic(self, mock_pm):
        """Test basic 2x2 grid layout creation."""
        mock_pm.get_presentation.return_value = self.pres
        
        elements = [
            {"content": "Item 1"},
            {"content": "Item 2"},
            {"content": "Item 3"},
            {"content": "Item 4"}
        ]
        
        result = self.engine.create_grid_layout(
            slide_index=0,
            elements=elements,
            rows=2,
            cols=2
        )
        
        self.assertNotIn("error", result)
        self.assertEqual(result["layout_type"], "grid")
        self.assertEqual(len(result["shapes"]), 4)
    
    @patch('layout_manager.presentation_manager')
    def test_create_grid_layout_invalid_slide(self, mock_pm):
        """Test grid layout with invalid slide index."""
        mock_pm.get_presentation.return_value = self.pres
        
        result = self.engine.create_grid_layout(
            slide_index=99,
            elements=[{"content": "Test"}],
            rows=1,
            cols=1
        )
        
        self.assertIn("error", result)
    
    @patch('layout_manager.presentation_manager')
    def test_create_grid_layout_3x3(self, mock_pm):
        """Test 3x3 grid layout creation."""
        mock_pm.get_presentation.return_value = self.pres
        
        elements = [{"content": f"Item {i}"} for i in range(9)]
        
        result = self.engine.create_grid_layout(
            slide_index=0,
            elements=elements,
            rows=3,
            cols=3
        )
        
        self.assertNotIn("error", result)
        self.assertEqual(len(result["shapes"]), 9)
    
    @patch('layout_manager.presentation_manager')
    def test_create_grid_layout_fewer_elements_than_cells(self, mock_pm):
        """Test grid layout with fewer elements than grid cells."""
        mock_pm.get_presentation.return_value = self.pres
        
        elements = [{"content": "Item 1"}, {"content": "Item 2"}]
        
        result = self.engine.create_grid_layout(
            slide_index=0,
            elements=elements,
            rows=2,
            cols=2
        )
        
        self.assertNotIn("error", result)
        self.assertEqual(len(result["shapes"]), 2)
    
    @patch('layout_manager.presentation_manager')
    def test_create_grid_layout_more_elements_than_cells(self, mock_pm):
        """Test grid layout with more elements than grid cells."""
        mock_pm.get_presentation.return_value = self.pres
        
        elements = [{"content": f"Item {i}"} for i in range(10)]
        
        result = self.engine.create_grid_layout(
            slide_index=0,
            elements=elements,
            rows=2,
            cols=2  # Only 4 cells
        )
        
        self.assertNotIn("error", result)
        # Should only create 4 shapes (max capacity)
        self.assertEqual(len(result["shapes"]), 4)
    
    @patch('layout_manager.presentation_manager')
    def test_create_grid_layout_with_custom_bounds(self, mock_pm):
        """Test grid layout with custom bounds."""
        mock_pm.get_presentation.return_value = self.pres
        
        custom_bounds = LayoutBounds(left=1.0, top=2.0, width=8.0, height=4.0)
        
        result = self.engine.create_grid_layout(
            slide_index=0,
            elements=[{"content": "Test"}],
            rows=1,
            cols=1,
            bounds=custom_bounds
        )
        
        self.assertNotIn("error", result)
    
    @patch('layout_manager.presentation_manager')
    def test_create_grid_layout_with_styled_elements(self, mock_pm):
        """Test grid layout with styled elements."""
        mock_pm.get_presentation.return_value = self.pres
        
        elements = [
            {"content": "Red", "fill_color": [255, 0, 0]},
            {"content": "Green", "fill_color": [0, 255, 0]}
        ]
        
        result = self.engine.create_grid_layout(
            slide_index=0,
            elements=elements,
            rows=1,
            cols=2
        )
        
        self.assertNotIn("error", result)


class TestCreateListLayout(unittest.TestCase):
    """Tests for create_list_layout method."""
    
    def setUp(self):
        """Set up test fixtures."""
        self.engine = LayoutEngine()
        self.pres = Presentation()
        self.pres.slides.add_slide(self.pres.slide_layouts[BLANK_SLIDE_LAYOUT_INDEX])
    
    @patch('layout_manager.presentation_manager')
    def test_create_vertical_list(self, mock_pm):
        """Test vertical list layout creation."""
        mock_pm.get_presentation.return_value = self.pres
        
        elements = [
            {"content": "Item 1"},
            {"content": "Item 2"},
            {"content": "Item 3"}
        ]
        
        result = self.engine.create_list_layout(
            slide_index=0,
            elements=elements,
            direction="vertical"
        )
        
        self.assertNotIn("error", result)
        self.assertEqual(result["layout_type"], "list")
        self.assertEqual(result["direction"], "vertical")
        self.assertEqual(len(result["shapes"]), 3)
    
    @patch('layout_manager.presentation_manager')
    def test_create_horizontal_list(self, mock_pm):
        """Test horizontal list layout creation."""
        mock_pm.get_presentation.return_value = self.pres
        
        elements = [{"content": f"Item {i}"} for i in range(4)]
        
        result = self.engine.create_list_layout(
            slide_index=0,
            elements=elements,
            direction="horizontal"
        )
        
        self.assertNotIn("error", result)
        self.assertEqual(result["direction"], "horizontal")
    
    @patch('layout_manager.presentation_manager')
    def test_create_list_layout_empty_elements(self, mock_pm):
        """Test list layout with no elements."""
        mock_pm.get_presentation.return_value = self.pres
        
        result = self.engine.create_list_layout(
            slide_index=0,
            elements=[],
            direction="vertical"
        )
        
        self.assertIn("error", result)
    
    @patch('layout_manager.presentation_manager')
    def test_create_list_layout_with_alignment(self, mock_pm):
        """Test list layout with different alignments."""
        mock_pm.get_presentation.return_value = self.pres
        
        elements = [{"content": "Item 1"}, {"content": "Item 2"}]
        
        # Test center alignment
        result = self.engine.create_list_layout(
            slide_index=0,
            elements=elements,
            direction="vertical",
            alignment="center"
        )
        
        self.assertNotIn("error", result)
    
    @patch('layout_manager.presentation_manager')
    def test_create_list_layout_invalid_slide(self, mock_pm):
        """Test list layout with invalid slide index."""
        mock_pm.get_presentation.return_value = self.pres
        
        result = self.engine.create_list_layout(
            slide_index=99,
            elements=[{"content": "Test"}],
            direction="vertical"
        )
        
        self.assertIn("error", result)


class TestCreateHierarchyLayout(unittest.TestCase):
    """Tests for create_hierarchy_layout method."""
    
    def setUp(self):
        """Set up test fixtures."""
        self.engine = LayoutEngine()
        self.pres = Presentation()
        self.pres.slides.add_slide(self.pres.slide_layouts[BLANK_SLIDE_LAYOUT_INDEX])
    
    @patch('layout_manager.presentation_manager')
    def test_create_simple_hierarchy(self, mock_pm):
        """Test simple hierarchy layout creation."""
        mock_pm.get_presentation.return_value = self.pres
        
        root = {
            "content": "Root",
            "children": [
                {"content": "Child 1"},
                {"content": "Child 2"}
            ]
        }
        
        result = self.engine.create_hierarchy_layout(
            slide_index=0,
            root=root
        )
        
        self.assertNotIn("error", result)
        self.assertEqual(result["layout_type"], "hierarchy")
        self.assertEqual(result["levels"], 2)
    
    @patch('layout_manager.presentation_manager')
    def test_create_deep_hierarchy(self, mock_pm):
        """Test deep hierarchy with multiple levels."""
        mock_pm.get_presentation.return_value = self.pres
        
        root = {
            "content": "CEO",
            "children": [
                {
                    "content": "VP",
                    "children": [
                        {"content": "Manager"}
                    ]
                }
            ]
        }
        
        result = self.engine.create_hierarchy_layout(
            slide_index=0,
            root=root
        )
        
        self.assertNotIn("error", result)
        self.assertEqual(result["levels"], 3)
    
    @patch('layout_manager.presentation_manager')
    def test_create_hierarchy_single_node(self, mock_pm):
        """Test hierarchy with single root node."""
        mock_pm.get_presentation.return_value = self.pres
        
        root = {"content": "Only Node"}
        
        result = self.engine.create_hierarchy_layout(
            slide_index=0,
            root=root
        )
        
        self.assertNotIn("error", result)
        self.assertEqual(result["levels"], 1)
    
    @patch('layout_manager.presentation_manager')
    def test_create_hierarchy_without_connectors(self, mock_pm):
        """Test hierarchy without connecting lines."""
        mock_pm.get_presentation.return_value = self.pres
        
        root = {
            "content": "Root",
            "children": [{"content": "Child"}]
        }
        
        result = self.engine.create_hierarchy_layout(
            slide_index=0,
            root=root,
            show_connectors=False
        )
        
        self.assertNotIn("error", result)
    
    @patch('layout_manager.presentation_manager')
    def test_create_hierarchy_invalid_slide(self, mock_pm):
        """Test hierarchy with invalid slide index."""
        mock_pm.get_presentation.return_value = self.pres
        
        result = self.engine.create_hierarchy_layout(
            slide_index=99,
            root={"content": "Test"}
        )
        
        self.assertIn("error", result)


class TestCreateFlowLayout(unittest.TestCase):
    """Tests for create_flow_layout method."""
    
    def setUp(self):
        """Set up test fixtures."""
        self.engine = LayoutEngine()
        self.pres = Presentation()
        self.pres.slides.add_slide(self.pres.slide_layouts[BLANK_SLIDE_LAYOUT_INDEX])
    
    @patch('layout_manager.presentation_manager')
    def test_create_horizontal_flow(self, mock_pm):
        """Test horizontal flow layout creation."""
        mock_pm.get_presentation.return_value = self.pres
        
        steps = [
            {"content": "Start"},
            {"content": "Process"},
            {"content": "End"}
        ]
        
        result = self.engine.create_flow_layout(
            slide_index=0,
            steps=steps,
            direction="horizontal"
        )
        
        self.assertNotIn("error", result)
        self.assertEqual(result["layout_type"], "flow")
        self.assertEqual(result["direction"], "horizontal")
        self.assertEqual(len(result["shapes"]), 3)
    
    @patch('layout_manager.presentation_manager')
    def test_create_vertical_flow(self, mock_pm):
        """Test vertical flow layout creation."""
        mock_pm.get_presentation.return_value = self.pres
        
        steps = [
            {"content": "Step 1"},
            {"content": "Step 2"}
        ]
        
        result = self.engine.create_flow_layout(
            slide_index=0,
            steps=steps,
            direction="vertical"
        )
        
        self.assertNotIn("error", result)
        self.assertEqual(result["direction"], "vertical")
    
    @patch('layout_manager.presentation_manager')
    def test_create_flow_empty_steps(self, mock_pm):
        """Test flow layout with no steps."""
        mock_pm.get_presentation.return_value = self.pres
        
        result = self.engine.create_flow_layout(
            slide_index=0,
            steps=[],
            direction="horizontal"
        )
        
        self.assertIn("error", result)
    
    @patch('layout_manager.presentation_manager')
    def test_create_flow_without_connectors(self, mock_pm):
        """Test flow layout without connectors."""
        mock_pm.get_presentation.return_value = self.pres
        
        steps = [{"content": "Step 1"}, {"content": "Step 2"}]
        
        result = self.engine.create_flow_layout(
            slide_index=0,
            steps=steps,
            show_connectors=False
        )
        
        self.assertNotIn("error", result)
    
    @patch('layout_manager.presentation_manager')
    def test_create_flow_with_line_connectors(self, mock_pm):
        """Test flow layout with line connectors instead of arrows."""
        mock_pm.get_presentation.return_value = self.pres
        
        steps = [{"content": "Step 1"}, {"content": "Step 2"}]
        
        result = self.engine.create_flow_layout(
            slide_index=0,
            steps=steps,
            connector_style="line"
        )
        
        self.assertNotIn("error", result)
    
    @patch('layout_manager.presentation_manager')
    def test_create_flow_invalid_slide(self, mock_pm):
        """Test flow layout with invalid slide index."""
        mock_pm.get_presentation.return_value = self.pres
        
        result = self.engine.create_flow_layout(
            slide_index=99,
            steps=[{"content": "Test"}]
        )
        
        self.assertIn("error", result)


class TestHelperMethods(unittest.TestCase):
    """Tests for helper methods in LayoutEngine."""
    
    def setUp(self):
        """Set up test fixtures."""
        self.engine = LayoutEngine()
    
    def test_create_element_from_dict_minimal(self):
        """Test creating element from minimal dict."""
        defaults = {
            "font_name": "Calibri",
            "font_size": 14,
            "fill_color": [79, 129, 189],
            "text_color": [0, 0, 0],
            "line_color": [0, 0, 0]
        }
        
        elem = self.engine._create_element_from_dict(
            {"content": "Test"},
            defaults
        )
        
        self.assertEqual(elem.content, "Test")
        self.assertEqual(elem.element_type, "textbox")
        self.assertEqual(elem.font_size, 14)
        self.assertEqual(elem.fill_color, [79, 129, 189])
    
    def test_create_element_from_dict_with_overrides(self):
        """Test creating element with override values."""
        defaults = {
            "font_name": "Calibri",
            "font_size": 14,
            "fill_color": [79, 129, 189],
            "text_color": [0, 0, 0],
            "line_color": [0, 0, 0]
        }
        
        elem = self.engine._create_element_from_dict(
            {
                "content": "Custom",
                "font_size": 24,
                "fill_color": [255, 0, 0],
                "bold": True
            },
            defaults
        )
        
        self.assertEqual(elem.font_size, 24)
        self.assertEqual(elem.fill_color, [255, 0, 0])
        self.assertTrue(elem.bold)
    
    @patch('layout_manager.template_manager')
    def test_get_default_styles(self, mock_tm):
        """Test getting default styles."""
        mock_tm.get_default_font_settings.return_value = {
            "body_font_name": "Arial",
            "body_font_size": 16
        }
        mock_tm.get_default_color_settings.return_value = {
            "accent_1": (100, 150, 200),
            "text_1": (50, 50, 50)
        }
        
        styles = self.engine._get_default_styles()
        
        self.assertEqual(styles["font_name"], "Arial")
        self.assertEqual(styles["font_size"], 16)
        self.assertEqual(styles["fill_color"], [100, 150, 200])
    
    @patch('layout_manager.template_manager')
    def test_resolve_color_with_semantic_tag(self, mock_tm):
        """Test color resolution with semantic tag."""
        mock_tm.resolve_color.return_value = [255, 128, 0]
        
        result = self.engine._resolve_color("accent")
        
        mock_tm.resolve_color.assert_called_once_with("accent")
        self.assertEqual(result, [255, 128, 0])
    
    @patch('layout_manager.template_manager')
    def test_resolve_color_with_rgb_list(self, mock_tm):
        """Test color resolution with RGB list."""
        mock_tm.resolve_color.return_value = [100, 200, 50]
        
        result = self.engine._resolve_color([100, 200, 50])
        
        self.assertEqual(result, [100, 200, 50])
    
    @patch('layout_manager.template_manager')
    def test_resolve_color_none(self, mock_tm):
        """Test color resolution with None input."""
        mock_tm.resolve_color.return_value = None
        
        result = self.engine._resolve_color(None)
        
        self.assertIsNone(result)


class TestFlattenHierarchy(unittest.TestCase):
    """Tests for _flatten_hierarchy method."""
    
    def setUp(self):
        """Set up test fixtures."""
        self.engine = LayoutEngine()
    
    def test_flatten_single_node(self):
        """Test flattening single node hierarchy."""
        root = {"content": "Root"}
        
        levels = self.engine._flatten_hierarchy(root)
        
        self.assertEqual(len(levels), 1)
        self.assertEqual(len(levels[0]), 1)
    
    def test_flatten_two_level_hierarchy(self):
        """Test flattening two-level hierarchy."""
        root = {
            "content": "Root",
            "children": [
                {"content": "Child 1"},
                {"content": "Child 2"}
            ]
        }
        
        levels = self.engine._flatten_hierarchy(root)
        
        self.assertEqual(len(levels), 2)
        self.assertEqual(len(levels[0]), 1)  # Root
        self.assertEqual(len(levels[1]), 2)  # Children
    
    def test_flatten_assigns_ids(self):
        """Test that flattening assigns unique IDs."""
        root = {
            "content": "Root",
            "children": [{"content": "Child"}]
        }
        
        levels = self.engine._flatten_hierarchy(root)
        
        # Check IDs are assigned
        self.assertIn("_id", levels[0][0])
        self.assertIn("_id", levels[1][0])
        self.assertNotEqual(levels[0][0]["_id"], levels[1][0]["_id"])
    
    def test_flatten_assigns_parent_ids(self):
        """Test that flattening assigns parent IDs."""
        root = {
            "content": "Root",
            "children": [{"content": "Child"}]
        }
        
        levels = self.engine._flatten_hierarchy(root)
        
        # Child should have parent ID
        self.assertIn("_parent_id", levels[1][0])
        self.assertEqual(levels[1][0]["_parent_id"], levels[0][0]["_id"])


class TestEnums(unittest.TestCase):
    """Tests for layout enumerations."""
    
    def test_layout_type_values(self):
        """Test LayoutType enum values."""
        self.assertEqual(LayoutType.GRID.value, "grid")
        self.assertEqual(LayoutType.LIST.value, "list")
        self.assertEqual(LayoutType.HIERARCHY.value, "hierarchy")
        self.assertEqual(LayoutType.FLOW.value, "flow")
    
    def test_alignment_values(self):
        """Test Alignment enum values."""
        self.assertEqual(Alignment.LEFT.value, "left")
        self.assertEqual(Alignment.CENTER.value, "center")
        self.assertEqual(Alignment.RIGHT.value, "right")
    
    def test_direction_values(self):
        """Test Direction enum values."""
        self.assertEqual(Direction.HORIZONTAL.value, "horizontal")
        self.assertEqual(Direction.VERTICAL.value, "vertical")
        self.assertEqual(Direction.LEFT_TO_RIGHT.value, "left_to_right")


if __name__ == '__main__':
    unittest.main()
