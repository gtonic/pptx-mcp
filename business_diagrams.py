"""
Business Diagrams Module

High-Level API for creating common business diagram types from structured data.
AI/LLM just provides the raw data (JSON/Dict), and this module automatically
generates professional layouts and visual arrangements.

Supported Diagrams:
- SWOT Analysis: 2x2 grid with Strengths, Weaknesses, Opportunities, Threats
- Timeline: Horizontal or vertical timeline with events/milestones
- Organization Chart: Hierarchical org structure with connectors

Goal: The AI only needs to provide the structured data (e.g., list of SWOT items),
not the visual arrangement. This module handles positioning, colors, and styling.
"""
from typing import Optional, Dict, Any, List, Union
from dataclasses import dataclass
import logging

from presentation_manager import presentation_manager
from layout_manager import layout_manager, LayoutBounds
from template_manager import template_manager
from performance_optimizer import performance_monitor
import ppt_utils

logger = logging.getLogger(__name__)


# ============================================================================
# DATA CLASSES FOR BUSINESS DIAGRAMS
# ============================================================================

@dataclass
class SWOTData:
    """Data structure for SWOT analysis."""
    strengths: List[str]
    weaknesses: List[str]
    opportunities: List[str]
    threats: List[str]
    title: Optional[str] = None


@dataclass
class TimelineEvent:
    """Single event on a timeline."""
    label: str
    date: Optional[str] = None
    description: Optional[str] = None
    color: Optional[Union[str, List[int]]] = None


@dataclass
class OrgChartNode:
    """Node in an organization chart."""
    name: str
    title: Optional[str] = None
    children: Optional[List['OrgChartNode']] = None
    color: Optional[Union[str, List[int]]] = None


# ============================================================================
# DEFAULT COLORS FOR BUSINESS DIAGRAMS
# ============================================================================

# SWOT default colors (semantic or RGB fallbacks)
SWOT_COLORS = {
    "strengths": "success",      # Green
    "weaknesses": "critical",    # Red  
    "opportunities": "info",     # Blue
    "threats": "warning"         # Yellow/Orange
}

SWOT_COLORS_RGB = {
    "strengths": [0, 176, 80],       # Green
    "weaknesses": [192, 0, 0],       # Red
    "opportunities": [0, 112, 192],  # Blue
    "threats": [255, 192, 0]         # Yellow/Orange
}

# Timeline default colors
TIMELINE_COLORS = {
    "connector": [128, 128, 128],    # Gray line
    "event_default": "primary",       # Primary theme color
    "milestone": "accent"             # Accent color for milestones
}


class BusinessDiagramsEngine:
    """
    High-Level Business Diagrams Engine.
    
    Provides specialized APIs for creating common business diagram types
    without requiring detailed coordinate specifications.
    """
    
    def _resolve_color(self, color_input: Union[str, List[int], None]) -> Optional[List[int]]:
        """Resolve color input to RGB values."""
        return template_manager.resolve_color(color_input)
    
    def _get_slide_bounds(self, presentation_id: Optional[str] = None) -> LayoutBounds:
        """Get layout bounds from presentation."""
        return layout_manager._get_slide_bounds(presentation_id)
    
    # =========================================================================
    # SWOT ANALYSIS
    # =========================================================================
    
    @performance_monitor.track_operation("create_swot_analysis")
    def create_swot_analysis(
        self,
        slide_index: int,
        strengths: List[str],
        weaknesses: List[str],
        opportunities: List[str],
        threats: List[str],
        title: Optional[str] = None,
        show_labels: bool = True,
        presentation_id: Optional[str] = None
    ) -> Dict[str, Any]:
        """
        Create a SWOT analysis diagram from structured data.
        
        Automatically generates a professional 2x2 grid layout with:
        - Top-left: Strengths (green)
        - Top-right: Weaknesses (red)
        - Bottom-left: Opportunities (blue)
        - Bottom-right: Threats (yellow/orange)
        
        Args:
            slide_index: Index of the target slide
            strengths: List of strength items
            weaknesses: List of weakness items
            opportunities: List of opportunity items
            threats: List of threat items
            title: Optional title for the SWOT diagram
            show_labels: Whether to show category labels (S, W, O, T)
            presentation_id: Optional presentation ID
            
        Returns:
            Dictionary with creation results and shape indices
            
        Example:
            create_swot_analysis(
                slide_index=0,
                strengths=["Strong brand", "Loyal customers", "Innovative products"],
                weaknesses=["High costs", "Limited distribution"],
                opportunities=["New markets", "Digital expansion"],
                threats=["Competition", "Economic downturn"]
            )
        """
        try:
            pres = presentation_manager.get_presentation(presentation_id)
            if not (0 <= slide_index < len(pres.slides)):
                return {"error": f"Invalid slide index: {slide_index}"}
            slide = pres.slides[slide_index]
            
            bounds = self._get_slide_bounds(presentation_id)
            
            # Adjust bounds if title is present
            if title:
                title_height = 0.6
                ppt_utils.add_textbox(
                    slide, bounds.left, bounds.top - 0.2, bounds.width, title_height, title,
                    font_size=24, bold=True, alignment="center"
                )
                bounds.top += title_height
                bounds.height -= title_height
            
            # Calculate grid dimensions
            gap = 0.15
            cell_width = (bounds.width - gap) / 2
            cell_height = (bounds.height - gap) / 2
            
            created_shapes = []
            
            # Define quadrants: (row, col, category, items, color_key)
            quadrants = [
                (0, 0, "Strengths", strengths, "strengths"),
                (0, 1, "Weaknesses", weaknesses, "weaknesses"),
                (1, 0, "Opportunities", opportunities, "opportunities"),
                (1, 1, "Threats", threats, "threats"),
            ]
            
            for row, col, label, items, color_key in quadrants:
                left = bounds.left + col * (cell_width + gap)
                top = bounds.top + row * (cell_height + gap)
                
                # Resolve color (try semantic first, then RGB fallback)
                fill_color = self._resolve_color(SWOT_COLORS[color_key])
                if fill_color is None:
                    fill_color = SWOT_COLORS_RGB[color_key]
                
                # Create the quadrant content
                if show_labels:
                    content = f"{label}\n\n" + "\n".join(f"• {item}" for item in items)
                else:
                    content = "\n".join(f"• {item}" for item in items)
                
                # Add shape with text
                shape = ppt_utils.add_shape(
                    slide, "rounded_rectangle",
                    left, top, cell_width, cell_height,
                    fill_color=fill_color,
                    line_color=[64, 64, 64]
                )
                
                # Add text to the shape
                if hasattr(shape, 'text_frame'):
                    from pptx.util import Pt
                    from pptx.dml.color import RGBColor
                    from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
                    
                    tf = shape.text_frame
                    tf.word_wrap = True
                    tf.auto_size = None
                    
                    # Clear and add content
                    p = tf.paragraphs[0]
                    p.text = content
                    p.alignment = PP_ALIGN.LEFT
                    p.font.size = Pt(11)
                    p.font.color.rgb = RGBColor(255, 255, 255)
                    
                    # If show_labels, make the first line bold
                    if show_labels and p.runs:
                        p.runs[0].font.bold = True
                
                created_shapes.append({
                    "index": len(slide.shapes) - 1,
                    "quadrant": label,
                    "bounds": {"left": left, "top": top, "width": cell_width, "height": cell_height}
                })
            
            return {
                "message": f"Created SWOT analysis with {sum(len(q[3]) for q in quadrants)} items",
                "slide_index": slide_index,
                "diagram_type": "swot_analysis",
                "shapes": created_shapes,
                "item_counts": {
                    "strengths": len(strengths),
                    "weaknesses": len(weaknesses),
                    "opportunities": len(opportunities),
                    "threats": len(threats)
                }
            }
            
        except (ValueError, KeyError) as e:
            return {"error": str(e)}
    
    # =========================================================================
    # TIMELINE
    # =========================================================================
    
    @performance_monitor.track_operation("create_timeline")
    def create_timeline(
        self,
        slide_index: int,
        events: List[Dict[str, Any]],
        direction: str = "horizontal",
        title: Optional[str] = None,
        show_connector: bool = True,
        presentation_id: Optional[str] = None
    ) -> Dict[str, Any]:
        """
        Create a timeline diagram from a list of events.
        
        Automatically generates a professional timeline with:
        - Central connector line
        - Event markers/nodes
        - Labels with dates and descriptions
        
        Args:
            slide_index: Index of the target slide
            events: List of event dictionaries, each with:
                - label (required): Event name/title
                - date (optional): Date or time period
                - description (optional): Additional details
                - color (optional): Semantic tag or RGB list
            direction: 'horizontal' (left to right) or 'vertical' (top to bottom)
            title: Optional title for the timeline
            show_connector: Whether to show the connecting line
            presentation_id: Optional presentation ID
            
        Returns:
            Dictionary with creation results and shape indices
            
        Example:
            create_timeline(
                slide_index=0,
                events=[
                    {"label": "Project Start", "date": "Jan 2024"},
                    {"label": "Phase 1 Complete", "date": "Mar 2024"},
                    {"label": "Launch", "date": "Jun 2024", "color": "success"}
                ],
                direction="horizontal"
            )
        """
        try:
            pres = presentation_manager.get_presentation(presentation_id)
            if not (0 <= slide_index < len(pres.slides)):
                return {"error": f"Invalid slide index: {slide_index}"}
            slide = pres.slides[slide_index]
            
            n_events = len(events)
            if n_events == 0:
                return {"error": "No events provided"}
            
            bounds = self._get_slide_bounds(presentation_id)
            
            # Adjust bounds if title is present
            if title:
                title_height = 0.6
                ppt_utils.add_textbox(
                    slide, bounds.left, bounds.top - 0.2, bounds.width, title_height, title,
                    font_size=24, bold=True, alignment="center"
                )
                bounds.top += title_height
                bounds.height -= title_height
            
            created_shapes = []
            is_horizontal = direction.lower() == "horizontal"
            
            if is_horizontal:
                # Horizontal timeline
                connector_y = bounds.top + bounds.height * 0.5
                event_spacing = bounds.width / (n_events + 1)
                node_size = 0.3
                
                # Draw connector line
                if show_connector:
                    ppt_utils.add_line(
                        slide,
                        bounds.left, connector_y,
                        bounds.left + bounds.width, connector_y,
                        line_color=TIMELINE_COLORS["connector"],
                        line_width=3.0
                    )
                    created_shapes.append({
                        "index": len(slide.shapes) - 1,
                        "type": "connector"
                    })
                
                # Add events
                for i, event in enumerate(events):
                    x = bounds.left + (i + 1) * event_spacing
                    
                    # Resolve event color
                    event_color = self._resolve_color(event.get("color", TIMELINE_COLORS["event_default"]))
                    if event_color is None:
                        event_color = [79, 129, 189]  # Default blue
                    
                    # Event node (circle)
                    ppt_utils.add_shape(
                        slide, "oval",
                        x - node_size/2, connector_y - node_size/2,
                        node_size, node_size,
                        fill_color=event_color,
                        line_color=[255, 255, 255],
                        line_width=2.0
                    )
                    node_shape_idx = len(slide.shapes) - 1
                    
                    # Determine label position (alternate above/below)
                    above = (i % 2 == 0)
                    label_y = connector_y - 0.8 if above else connector_y + 0.5
                    
                    # Build label text
                    label_text = event.get("label", "")
                    date_text = event.get("date", "")
                    desc_text = event.get("description", "")
                    
                    full_text = label_text
                    if date_text:
                        full_text = f"{date_text}\n{full_text}"
                    if desc_text:
                        full_text = f"{full_text}\n{desc_text}"
                    
                    # Event label
                    label_width = event_spacing * 0.9
                    label_height = 0.6
                    ppt_utils.add_textbox(
                        slide,
                        x - label_width/2, label_y,
                        label_width, label_height,
                        full_text,
                        font_size=10,
                        alignment="center"
                    )
                    label_shape_idx = len(slide.shapes) - 1
                    
                    created_shapes.append({
                        "index": node_shape_idx,
                        "type": "event_node",
                        "event_index": i,
                        "label": event.get("label", "")
                    })
                    created_shapes.append({
                        "index": label_shape_idx,
                        "type": "event_label",
                        "event_index": i
                    })
                    
            else:
                # Vertical timeline
                connector_x = bounds.left + bounds.width * 0.3
                event_spacing = bounds.height / (n_events + 1)
                node_size = 0.3
                
                # Draw connector line
                if show_connector:
                    ppt_utils.add_line(
                        slide,
                        connector_x, bounds.top,
                        connector_x, bounds.top + bounds.height,
                        line_color=TIMELINE_COLORS["connector"],
                        line_width=3.0
                    )
                    created_shapes.append({
                        "index": len(slide.shapes) - 1,
                        "type": "connector"
                    })
                
                # Add events
                for i, event in enumerate(events):
                    y = bounds.top + (i + 1) * event_spacing
                    
                    # Resolve event color
                    event_color = self._resolve_color(event.get("color", TIMELINE_COLORS["event_default"]))
                    if event_color is None:
                        event_color = [79, 129, 189]
                    
                    # Event node (circle)
                    ppt_utils.add_shape(
                        slide, "oval",
                        connector_x - node_size/2, y - node_size/2,
                        node_size, node_size,
                        fill_color=event_color,
                        line_color=[255, 255, 255],
                        line_width=2.0
                    )
                    node_shape_idx = len(slide.shapes) - 1
                    
                    # Build label text
                    label_text = event.get("label", "")
                    date_text = event.get("date", "")
                    desc_text = event.get("description", "")
                    
                    full_text = label_text
                    if date_text:
                        full_text = f"{date_text}: {full_text}"
                    if desc_text:
                        full_text = f"{full_text}\n{desc_text}"
                    
                    # Event label (to the right of the node)
                    label_left = connector_x + node_size
                    label_width = bounds.width - (connector_x - bounds.left) - node_size - 0.2
                    label_height = 0.5
                    
                    ppt_utils.add_textbox(
                        slide,
                        label_left, y - label_height/2,
                        label_width, label_height,
                        full_text,
                        font_size=11,
                        alignment="left"
                    )
                    label_shape_idx = len(slide.shapes) - 1
                    
                    created_shapes.append({
                        "index": node_shape_idx,
                        "type": "event_node",
                        "event_index": i,
                        "label": event.get("label", "")
                    })
                    created_shapes.append({
                        "index": label_shape_idx,
                        "type": "event_label",
                        "event_index": i
                    })
            
            return {
                "message": f"Created {direction} timeline with {n_events} events",
                "slide_index": slide_index,
                "diagram_type": "timeline",
                "direction": direction,
                "event_count": n_events,
                "shapes": created_shapes
            }
            
        except (ValueError, KeyError) as e:
            return {"error": str(e)}
    
    # =========================================================================
    # ORGANIZATION CHART
    # =========================================================================
    
    @performance_monitor.track_operation("create_org_chart")
    def create_org_chart(
        self,
        slide_index: int,
        root: Dict[str, Any],
        title: Optional[str] = None,
        show_connectors: bool = True,
        compact: bool = False,
        presentation_id: Optional[str] = None
    ) -> Dict[str, Any]:
        """
        Create an organization chart from hierarchical data.
        
        Automatically generates a professional org chart with:
        - Rectangular boxes for each person/role
        - Connecting lines between hierarchical levels
        - Automatic positioning based on tree structure
        
        Args:
            slide_index: Index of the target slide
            root: Root node of the organization, structured as:
                {
                    "name": "Person Name",
                    "title": "Job Title (optional)",
                    "children": [...] (optional list of child nodes with same structure),
                    "color": semantic tag or RGB list (optional)
                }
            title: Optional title for the org chart
            show_connectors: Whether to show connecting lines
            compact: Whether to use compact layout (smaller boxes)
            presentation_id: Optional presentation ID
            
        Returns:
            Dictionary with creation results and shape indices
            
        Example:
            create_org_chart(
                slide_index=0,
                root={
                    "name": "John Smith",
                    "title": "CEO",
                    "children": [
                        {
                            "name": "Jane Doe",
                            "title": "VP Sales",
                            "children": [
                                {"name": "Alice", "title": "Sales Lead"},
                                {"name": "Bob", "title": "Account Manager"}
                            ]
                        },
                        {
                            "name": "Mike Wilson",
                            "title": "VP Engineering",
                            "children": [
                                {"name": "Carol", "title": "Tech Lead"}
                            ]
                        }
                    ]
                }
            )
        """
        try:
            pres = presentation_manager.get_presentation(presentation_id)
            if not (0 <= slide_index < len(pres.slides)):
                return {"error": f"Invalid slide index: {slide_index}"}
            
            bounds = self._get_slide_bounds(presentation_id)
            
            # Adjust bounds if title is present
            if title:
                slide = pres.slides[slide_index]
                title_height = 0.6
                ppt_utils.add_textbox(
                    slide, bounds.left, bounds.top - 0.2, bounds.width, title_height, title,
                    font_size=24, bold=True, alignment="center"
                )
                # Create new bounds with adjusted top
                bounds = LayoutBounds(
                    left=bounds.left,
                    top=bounds.top + title_height,
                    width=bounds.width,
                    height=bounds.height - title_height
                )
            
            # Transform the org chart data structure for the hierarchy layout
            hierarchy_root = self._transform_org_node(root, compact)
            
            # Use the existing hierarchy layout from layout_manager
            level_gap = 0.6 if compact else 0.8
            sibling_gap = 0.2 if compact else 0.3
            
            result = layout_manager.create_hierarchy_layout(
                slide_index=slide_index,
                root=hierarchy_root,
                level_gap=level_gap,
                sibling_gap=sibling_gap,
                show_connectors=show_connectors,
                bounds=bounds,
                presentation_id=presentation_id
            )
            
            # Add org chart specific metadata
            result["diagram_type"] = "org_chart"
            result["message"] = f"Created organization chart with {result.get('levels', 0)} levels"
            
            return result
            
        except (ValueError, KeyError) as e:
            return {"error": str(e)}
    
    def _transform_org_node(self, node: Dict[str, Any], compact: bool = False) -> Dict[str, Any]:
        """Transform org chart node to hierarchy layout format."""
        name = node.get("name", "")
        title = node.get("title", "")
        
        # Build content: Name on first line, title on second
        if title:
            content = f"{name}\n{title}"
        else:
            content = name
        
        # Resolve color
        color = node.get("color")
        fill_color = self._resolve_color(color) if color else None
        if fill_color is None:
            fill_color = [79, 129, 189]  # Default blue
        
        result = {
            "content": content,
            "element_type": "shape",
            "shape_type": "rounded_rectangle",
            "fill_color": fill_color,
            "text_color": [255, 255, 255],
            "font_size": 10 if compact else 12,
            "bold": True
        }
        
        # Recursively transform children
        children = node.get("children", [])
        if children:
            result["children"] = [
                self._transform_org_node(child, compact)
                for child in children
            ]
        
        return result


# Global instance
business_diagrams = BusinessDiagramsEngine()
