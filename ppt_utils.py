# ppt_utils.py
from typing import Optional, Dict, Any, List, Tuple
import base64
from io import BytesIO

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.chart.data import CategoryChartData
from pptx.enum.chart import XL_CHART_TYPE, XL_LEGEND_POSITION

DATA_DIR = "/data"

def create_presentation() -> Presentation:
    """Creates a new PowerPoint presentation object."""
    return Presentation()

import os

def open_presentation(file_path: str) -> Presentation:
    """Opens an existing PowerPoint presentation from a file."""
    # Ensure file_path is in DATA_DIR
    if not file_path.startswith(DATA_DIR + "/"):
        file_path = os.path.join(DATA_DIR, os.path.basename(file_path))
    return Presentation(file_path)

def save_presentation(pres: Presentation, file_path: str) -> str:
    """Saves a presentation object to a file."""
    # Ensure file_path is in DATA_DIR
    import os
    if not file_path.startswith(DATA_DIR + "/"):
        file_path = os.path.join(DATA_DIR, os.path.basename(file_path))
    pres.save(file_path)
    return file_path

def get_presentation_info(pres: Presentation) -> Dict:
    """Gets information about a presentation."""
    layouts = {i: layout.name for i, layout in enumerate(pres.slide_layouts)}
    props = pres.core_properties
    core_props = {
        "title": props.title, "subject": props.subject, "author": props.author,
        "keywords": props.keywords, "comments": props.comments
    }
    return {"slide_count": len(pres.slides), "slide_layouts": layouts, "core_properties": core_props}

def set_core_properties(pres: Presentation, **kwargs) -> Dict:
    """Sets core document properties."""
    props = pres.core_properties
    for key, value in kwargs.items():
        if value is not None:
            setattr(props, key, value)
    return get_presentation_info(pres)['core_properties']

def add_slide(pres: Presentation, layout_index: int, title: Optional[str] = None) -> Tuple[Any, Dict]:
    """Adds a new slide to the presentation."""
    slide_layout = pres.slide_layouts[layout_index]
    slide = pres.slides.add_slide(slide_layout)
    if title and slide.shapes.title:
        slide.shapes.title.text = title
    placeholders = {p.placeholder_format.idx: p.name for p in slide.placeholders}
    return slide, {"layout_name": slide_layout.name, "placeholders": placeholders}

def add_textbox(slide: Any, left: float, top: float, width: float, height: float, text: str, **kwargs):
    """Adds a textbox to a slide and formats it."""
    from pptx.util import Inches, Pt
    from pptx.dml.color import RGBColor
    from pptx.enum.text import PP_ALIGN
    txBox = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = txBox.text_frame
    tf.text = text
    p = tf.paragraphs[0]
    font = p.font
    if kwargs.get('font_size'): font.size = Pt(kwargs['font_size'])
    if kwargs.get('font_name'): font.name = kwargs['font_name']
    if kwargs.get('bold') is not None: font.bold = kwargs['bold']
    if kwargs.get('italic') is not None: font.italic = kwargs['italic']
    if kwargs.get('color'): font.color.rgb = RGBColor(*kwargs['color'])
    if kwargs.get('alignment'): p.alignment = getattr(PP_ALIGN, kwargs['alignment'].upper())
    return txBox

def add_shape(slide: Any, shape_type: str, left: float, top: float, width: float, height: float, **kwargs):
    """Adds an auto shape to a slide and formats it."""
    from pptx.util import Inches, Pt
    from pptx.dml.color import RGBColor
    shape_type_map = {
        'rectangle': 1, 'rounded_rectangle': 2, 'oval': 9, 'diamond': 4,
        'triangle': 5, 'right_triangle': 6, 'pentagon': 56, 'hexagon': 10,
        'heptagon': 11, 'octagon': 12, 'star': 12, 'arrow': 13, 'cloud': 35,
        'heart': 21, 'lightning_bolt': 22, 'sun': 23, 'moon': 24,
        'smiley_face': 17, 'no_symbol': 19, 'flowchart_process': 112,
        'flowchart_decision': 114, 'flowchart_data': 115, 'flowchart_document': 119
    }
    shape_type_lower = str(shape_type).lower()
    if shape_type_lower not in shape_type_map:
        raise ValueError(f"Unsupported shape type: '{shape_type}'")
    
    shape_value = shape_type_map[shape_type_lower]
    shape = slide.shapes.add_shape(shape_value, Inches(left), Inches(top), Inches(width), Inches(height))

    if kwargs.get('fill_color'):
        fill = shape.fill
        fill.solid()
        fill.fore_color.rgb = RGBColor(*kwargs['fill_color'])
    if kwargs.get('line_color'):
        line = shape.line
        line.color.rgb = RGBColor(*kwargs['line_color'])
        if kwargs.get('line_width'):
            line.width = Pt(kwargs['line_width'])
    return shape

def add_line(slide: Any, x1: float, y1: float, x2: float, y2: float, line_color: Optional[List[int]] = None, line_width: Optional[float] = None):
    """
    Adds a straight line (connector) to a slide.
    Coordinates are in inches.
    """
    from pptx.util import Inches, Pt
    from pptx.dml.color import RGBColor
    from pptx.enum.shapes import MSO_CONNECTOR

    line = slide.shapes.add_connector(
        MSO_CONNECTOR.STRAIGHT,
        Inches(x1), Inches(y1), Inches(x2), Inches(y2)
    )
    if line_color:
        line.line.color.rgb = RGBColor(*line_color)
    if line_width:
        line.line.width = Pt(line_width)
    return line

def extract_template_styles(pres: Presentation) -> Dict[str, Any]:
    """
    Extracts theme colors and default font info from a presentation.
    Returns a dict with 'colors' and 'fonts' keys.
    """
    # Extract theme colors
    colors = {}
    try:
        color_scheme = pres.theme_color_scheme
        for color_name in color_scheme.__dir__():
            if not color_name.startswith("_") and color_name not in ["part", "element"]:
                try:
                    color = getattr(color_scheme, color_name)
                    if hasattr(color, "rgb"):
                        colors[color_name] = tuple(color.rgb)
                except Exception:
                    continue
    except Exception:
        pass

    # Extract default font info from slide master
    fonts = {}
    try:
        master = pres.slide_master
        # Title and body font
        title_font = master.placeholders.title.text_frame.paragraphs[0].font
        body_font = master.placeholders[1].text_frame.paragraphs[0].font
        fonts["title_font_name"] = title_font.name
        fonts["title_font_size"] = title_font.size.pt if title_font.size else None
        fonts["body_font_name"] = body_font.name
        fonts["body_font_size"] = body_font.size.pt if body_font.size else None
    except Exception:
        pass

    return {"colors": colors, "fonts": fonts}
