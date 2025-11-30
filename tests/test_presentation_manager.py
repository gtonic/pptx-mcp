#!/usr/bin/env python3
"""
Tests for the presentation manager module.

Tests the core presentation lifecycle including creation, opening, saving,
and managing presentation state.
"""

import sys
import os
sys.path.insert(0, os.path.dirname(os.path.dirname(os.path.abspath(__file__))))

import unittest
from unittest.mock import patch, MagicMock
from pptx import Presentation

from presentation_manager import PresentationManager, presentation_manager


class TestPresentationManagerInitialization(unittest.TestCase):
    """Tests for PresentationManager initialization."""
    
    def test_initialization(self):
        """Test that PresentationManager initializes correctly."""
        pm = PresentationManager()
        self.assertEqual(pm.presentations, {})
        self.assertIsNone(pm.current_presentation_id)
    
    def test_global_instance_exists(self):
        """Test that global presentation_manager instance exists."""
        self.assertIsInstance(presentation_manager, PresentationManager)


class TestCreatePresentation(unittest.TestCase):
    """Tests for create_presentation method."""
    
    def setUp(self):
        """Set up test fixtures."""
        self.pm = PresentationManager()
    
    @patch('presentation_manager.ppt_utils.create_presentation')
    def test_create_presentation_with_auto_id(self, mock_create):
        """Test creating a presentation with auto-generated ID."""
        mock_pres = MagicMock()
        mock_pres.slides = []
        mock_create.return_value = mock_pres
        
        result = self.pm.create_presentation()
        
        self.assertIn("presentation_id", result)
        self.assertEqual(result["presentation_id"], "presentation_1")
        self.assertIn("message", result)
        self.assertEqual(result["slide_count"], 0)
        mock_create.assert_called_once()
    
    @patch('presentation_manager.ppt_utils.create_presentation')
    def test_create_presentation_with_custom_id(self, mock_create):
        """Test creating a presentation with custom ID."""
        mock_pres = MagicMock()
        mock_pres.slides = []
        mock_create.return_value = mock_pres
        
        result = self.pm.create_presentation(id="my_presentation")
        
        self.assertEqual(result["presentation_id"], "my_presentation")
        self.assertIn("my_presentation", self.pm.presentations)
    
    @patch('presentation_manager.ppt_utils.create_presentation')
    def test_create_presentation_sets_current(self, mock_create):
        """Test that create_presentation sets current presentation."""
        mock_pres = MagicMock()
        mock_pres.slides = []
        mock_create.return_value = mock_pres
        
        self.pm.create_presentation(id="test_pres")
        
        self.assertEqual(self.pm.current_presentation_id, "test_pres")
    
    @patch('presentation_manager.ppt_utils.create_presentation')
    def test_create_multiple_presentations(self, mock_create):
        """Test creating multiple presentations."""
        mock_pres = MagicMock()
        mock_pres.slides = []
        mock_create.return_value = mock_pres
        
        result1 = self.pm.create_presentation()
        result2 = self.pm.create_presentation()
        
        self.assertEqual(result1["presentation_id"], "presentation_1")
        self.assertEqual(result2["presentation_id"], "presentation_2")
        self.assertEqual(len(self.pm.presentations), 2)


class TestGetCurrentPresentation(unittest.TestCase):
    """Tests for get_current_presentation method."""
    
    def setUp(self):
        """Set up test fixtures."""
        self.pm = PresentationManager()
    
    def test_get_current_presentation_raises_when_none(self):
        """Test that getting current presentation raises when none exists."""
        with self.assertRaises(ValueError) as context:
            self.pm.get_current_presentation()
        
        self.assertIn("No presentation is currently loaded", str(context.exception))
    
    @patch('presentation_manager.ppt_utils.create_presentation')
    def test_get_current_presentation_returns_correct(self, mock_create):
        """Test that get_current_presentation returns the correct presentation."""
        mock_pres = MagicMock()
        mock_pres.slides = []
        mock_create.return_value = mock_pres
        
        self.pm.create_presentation()
        result = self.pm.get_current_presentation()
        
        self.assertEqual(result, mock_pres)
    
    def test_get_current_presentation_raises_when_id_not_found(self):
        """Test that get_current_presentation raises when ID not in presentations."""
        self.pm.current_presentation_id = "nonexistent_id"
        
        with self.assertRaises(ValueError):
            self.pm.get_current_presentation()


class TestOpenPresentation(unittest.TestCase):
    """Tests for open_presentation method."""
    
    def setUp(self):
        """Set up test fixtures."""
        self.pm = PresentationManager()
    
    @patch('presentation_manager.os.path.exists')
    def test_open_presentation_file_not_found(self, mock_exists):
        """Test opening a non-existent presentation."""
        mock_exists.return_value = False
        
        result = self.pm.open_presentation("/data/nonexistent.pptx")
        
        self.assertIn("error", result)
        self.assertIn("not found", result["error"])
    
    @patch('presentation_manager.ppt_utils.open_presentation')
    @patch('presentation_manager.os.path.exists')
    def test_open_presentation_success(self, mock_exists, mock_open):
        """Test successfully opening a presentation."""
        mock_exists.return_value = True
        mock_pres = MagicMock()
        mock_pres.slides = [MagicMock(), MagicMock()]
        mock_open.return_value = mock_pres
        
        result = self.pm.open_presentation("/data/test.pptx")
        
        self.assertNotIn("error", result)
        self.assertIn("presentation_id", result)
        self.assertEqual(result["slide_count"], 2)
    
    @patch('presentation_manager.ppt_utils.open_presentation')
    @patch('presentation_manager.os.path.exists')
    def test_open_presentation_with_custom_id(self, mock_exists, mock_open):
        """Test opening a presentation with custom ID."""
        mock_exists.return_value = True
        mock_pres = MagicMock()
        mock_pres.slides = []
        mock_open.return_value = mock_pres
        
        result = self.pm.open_presentation("/data/test.pptx", id="custom_id")
        
        self.assertEqual(result["presentation_id"], "custom_id")
    
    @patch('presentation_manager.ppt_utils.open_presentation')
    @patch('presentation_manager.os.path.exists')
    def test_open_presentation_handles_exception(self, mock_exists, mock_open):
        """Test opening a presentation handles exceptions."""
        mock_exists.return_value = True
        mock_open.side_effect = Exception("Failed to read file")
        
        result = self.pm.open_presentation("/data/test.pptx")
        
        self.assertIn("error", result)
        self.assertIn("Failed to open presentation", result["error"])
    
    @patch('presentation_manager.ppt_utils.open_presentation')
    @patch('presentation_manager.os.path.exists')
    def test_open_presentation_normalizes_path(self, mock_exists, mock_open):
        """Test that open_presentation normalizes file paths."""
        mock_exists.return_value = True
        mock_pres = MagicMock()
        mock_pres.slides = []
        mock_open.return_value = mock_pres
        
        # Pass a path without /data/ prefix
        result = self.pm.open_presentation("test.pptx")
        
        # Should normalize to /data/test.pptx
        mock_exists.assert_called()
        call_args = mock_exists.call_args[0][0]
        self.assertIn("/data/", call_args)


class TestSavePresentation(unittest.TestCase):
    """Tests for save_presentation method."""
    
    def setUp(self):
        """Set up test fixtures."""
        self.pm = PresentationManager()
    
    def test_save_presentation_no_current(self):
        """Test saving when no presentation is loaded."""
        result = self.pm.save_presentation("/data/test.pptx")
        
        self.assertIn("error", result)
    
    @patch('presentation_manager.ppt_utils.save_presentation')
    @patch('presentation_manager.ppt_utils.create_presentation')
    def test_save_presentation_success(self, mock_create, mock_save):
        """Test successfully saving a presentation."""
        mock_pres = MagicMock()
        mock_pres.slides = []
        mock_create.return_value = mock_pres
        mock_save.return_value = "/data/test.pptx"
        
        self.pm.create_presentation()
        result = self.pm.save_presentation("/data/test.pptx")
        
        self.assertNotIn("error", result)
        self.assertIn("message", result)
        self.assertEqual(result["file_path"], "/data/test.pptx")
    
    @patch('presentation_manager.ppt_utils.save_presentation')
    @patch('presentation_manager.ppt_utils.create_presentation')
    def test_save_presentation_by_id(self, mock_create, mock_save):
        """Test saving a presentation by ID."""
        mock_pres = MagicMock()
        mock_pres.slides = []
        mock_create.return_value = mock_pres
        mock_save.return_value = "/data/test.pptx"
        
        self.pm.create_presentation(id="my_pres")
        result = self.pm.save_presentation("/data/test.pptx", presentation_id="my_pres")
        
        self.assertNotIn("error", result)


class TestGetPresentationInfo(unittest.TestCase):
    """Tests for get_presentation_info method."""
    
    def setUp(self):
        """Set up test fixtures."""
        self.pm = PresentationManager()
    
    def test_get_presentation_info_no_current(self):
        """Test getting info when no presentation is loaded."""
        result = self.pm.get_presentation_info()
        
        self.assertIn("error", result)
    
    @patch('presentation_manager.ppt_utils.get_presentation_info')
    @patch('presentation_manager.ppt_utils.create_presentation')
    def test_get_presentation_info_success(self, mock_create, mock_info):
        """Test successfully getting presentation info."""
        mock_pres = MagicMock()
        mock_pres.slides = []
        mock_create.return_value = mock_pres
        mock_info.return_value = {
            "slide_count": 0,
            "slide_layouts": {0: "Title Slide"},
            "core_properties": {"title": "Test"}
        }
        
        self.pm.create_presentation(id="test_pres")
        result = self.pm.get_presentation_info()
        
        self.assertNotIn("error", result)
        self.assertEqual(result["presentation_id"], "test_pres")
        self.assertEqual(result["slide_count"], 0)


class TestSetCoreProperties(unittest.TestCase):
    """Tests for set_core_properties method."""
    
    def setUp(self):
        """Set up test fixtures."""
        self.pm = PresentationManager()
    
    def test_set_core_properties_no_current(self):
        """Test setting properties when no presentation is loaded."""
        result = self.pm.set_core_properties(title="Test")
        
        self.assertIn("error", result)
    
    @patch('presentation_manager.ppt_utils.set_core_properties')
    @patch('presentation_manager.ppt_utils.create_presentation')
    def test_set_core_properties_success(self, mock_create, mock_set_props):
        """Test successfully setting core properties."""
        mock_pres = MagicMock()
        mock_pres.slides = []
        mock_create.return_value = mock_pres
        mock_set_props.return_value = {
            "title": "My Title",
            "author": "Test Author"
        }
        
        self.pm.create_presentation()
        result = self.pm.set_core_properties(title="My Title", author="Test Author")
        
        self.assertNotIn("error", result)
        self.assertIn("message", result)
        self.assertIn("core_properties", result)


class TestGetPresentation(unittest.TestCase):
    """Tests for get_presentation method."""
    
    def setUp(self):
        """Set up test fixtures."""
        self.pm = PresentationManager()
    
    def test_get_presentation_no_id_no_current(self):
        """Test get_presentation without ID when no current exists."""
        with self.assertRaises(ValueError):
            self.pm.get_presentation()
    
    @patch('presentation_manager.ppt_utils.create_presentation')
    def test_get_presentation_returns_current(self, mock_create):
        """Test get_presentation returns current when no ID provided."""
        mock_pres = MagicMock()
        mock_pres.slides = []
        mock_create.return_value = mock_pres
        
        self.pm.create_presentation()
        result = self.pm.get_presentation()
        
        self.assertEqual(result, mock_pres)
    
    @patch('presentation_manager.ppt_utils.create_presentation')
    def test_get_presentation_by_id(self, mock_create):
        """Test get_presentation with specific ID."""
        mock_pres1 = MagicMock()
        mock_pres1.slides = []
        mock_pres2 = MagicMock()
        mock_pres2.slides = []
        mock_create.side_effect = [mock_pres1, mock_pres2]
        
        self.pm.create_presentation(id="pres1")
        self.pm.create_presentation(id="pres2")
        
        result = self.pm.get_presentation("pres1")
        
        self.assertEqual(result, mock_pres1)
    
    def test_get_presentation_nonexistent_id(self):
        """Test get_presentation with nonexistent ID."""
        with self.assertRaises(KeyError) as context:
            self.pm.get_presentation("nonexistent")
        
        self.assertIn("nonexistent", str(context.exception))


if __name__ == '__main__':
    unittest.main()
