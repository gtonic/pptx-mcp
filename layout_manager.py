"""
High-Level Layout Engine Module

Provides automatic layout calculation and placement for AI-generated PowerPoint diagrams.
This module eliminates the need for explicit pixel coordinates by offering structural
layout descriptions (grid, list, hierarchy, flow) that automatically calculate positions.
"""
from typing import Optional, Dict, Any, List, Union, Tuple
from dataclasses import dataclass, field
from enum import Enum
import math

from presentation_manager import presentation_manager
from template_manager import template_manager
from performance_optimizer import performance_monitor
import ppt_utils


class LayoutType(str, Enum):
    """Enumeration of available layout types."""
    GRID = "grid"
    LIST = "list"
    HIERARCHY = "hierarchy"
    FLOW = "flow"


class Alignment(str, Enum):
    """Alignment options for layout elements."""
    LEFT = "left"
    CENTER = "center"
    RIGHT = "right"
    TOP = "top"
    MIDDLE = "middle"
    BOTTOM = "bottom"


class Direction(str, Enum):
    """Direction for list and flow layouts."""
    HORIZONTAL = "horizontal"
    VERTICAL = "vertical"
    LEFT_TO_RIGHT = "left_to_right"
    RIGHT_TO_LEFT = "right_to_left"
    TOP_TO_BOTTOM = "top_to_bottom"
    BOTTOM_TO_TOP = "bottom_to_top"


@dataclass
class LayoutBounds:
    """Represents the bounding area for a layout."""
    left: float = 0.5  # inches from left edge
    top: float = 1.0   # inches from top edge (leaving room for title)
    width: float = 9.0  # inches
    height: float = 5.5  # inches

    @classmethod
    def from_slide(cls, slide_width: float = 10.0, slide_height: float = 7.5,
                   margin: float = 0.5, title_height: float = 1.0) -> 'LayoutBounds':
        """Create bounds from slide dimensions with margins."""
        return cls(
            left=margin,
            top=margin + title_height,
            width=slide_width - (2 * margin),
            height=slide_height - (2 * margin) - title_height
        )


@dataclass 
class LayoutElement:
    """Represents a single element in a layout."""
    content: str
    element_type: str = "textbox"  # textbox, shape, or image
    shape_type: Optional[str] = None  # For shape elements: rectangle, oval, etc.
    fill_color: Optional[List[int]] = None
    line_color: Optional[List[int]] = None
    font_size: Optional[int] = None
    font_name: Optional[str] = None
    bold: Optional[bool] = None
    text_color: Optional[List[int]] = None
    alignment: Optional[str] = "center"
    # Computed position (set by layout engine)
    left: float = 0.0
    top: float = 0.0
    width: float = 1.0
    height: float = 1.0


@dataclass
class GridLayoutConfig:
    """Configuration for grid layout."""
    rows: int = 2
    cols: int = 2
    gap: float = 0.2  # Gap between cells in inches
    equal_size: bool = True  # Whether all cells should be equal size


@dataclass
class ListLayoutConfig:
    """Configuration for list layout."""
    direction: Direction = Direction.VERTICAL
    gap: float = 0.2  # Gap between items in inches
    alignment: Alignment = Alignment.LEFT
    numbered: bool = False
    bullet_style: Optional[str] = None  # None, "circle", "square", "arrow"


@dataclass
class HierarchyLayoutConfig:
    """Configuration for hierarchy/tree layout."""
    direction: Direction = Direction.TOP_TO_BOTTOM
    level_gap: float = 0.8  # Vertical gap between levels in inches
    sibling_gap: float = 0.3  # Horizontal gap between siblings in inches
    connector_style: str = "line"  # "line", "elbow", "none"


@dataclass
class FlowLayoutConfig:
    """Configuration for flow/process layout."""
    direction: Direction = Direction.LEFT_TO_RIGHT
    gap: float = 0.4  # Gap between elements in inches
    connector_style: str = "arrow"  # "arrow", "line", "none"
    shape_type: str = "rounded_rectangle"


class LayoutEngine:
    """
    High-Level Layout Engine for automatic element positioning.
    
    This engine calculates positions and sizes based on structural descriptions,
    eliminating the need for AI/LLMs to provide explicit coordinates.
    """
    
    def __init__(self):
        self.default_bounds = LayoutBounds()
    
    def _get_slide_bounds(self, presentation_id: Optional[str] = None) -> LayoutBounds:
        """Get the layout bounds based on current presentation slide size."""
        try:
            pres = presentation_manager.get_presentation(presentation_id)
            return LayoutBounds.from_slide(
                slide_width=pres.slide_width.inches,
                slide_height=pres.slide_height.inches
            )
        except (ValueError, KeyError):
            return self.default_bounds
    
    def _get_default_styles(self) -> Dict[str, Any]:
        """Get default styling from template or sensible defaults."""
        fonts = template_manager.get_default_font_settings()
        colors = template_manager.get_default_color_settings()
        return {
            "font_name": fonts.get("body_font_name", "Calibri"),
            "font_size": fonts.get("body_font_size", 14),
            "fill_color": list(colors.get("accent_1", (79, 129, 189))),
            "text_color": list(colors.get("text_1", (0, 0, 0))),
            "line_color": list(colors.get("text_1", (0, 0, 0)))
        }
    
    # =========================================================================
    # GRID LAYOUT
    # =========================================================================
    
    @performance_monitor.track_operation("layout_grid")
    def create_grid_layout(
        self,
        slide_index: int,
        elements: List[Dict[str, Any]],
        rows: int = 2,
        cols: int = 2,
        gap: float = 0.2,
        bounds: Optional[LayoutBounds] = None,
        presentation_id: Optional[str] = None
    ) -> Dict[str, Any]:
        """
        Create elements arranged in a grid layout.
        
        Args:
            slide_index: Index of the target slide
            elements: List of element dictionaries with 'content' and optional styling
            rows: Number of rows in the grid
            cols: Number of columns in the grid  
            gap: Gap between cells in inches
            bounds: Optional custom bounds for the layout area
            presentation_id: Optional presentation ID
            
        Returns:
            Dictionary with confirmation and created shape indices
        """
        try:
            pres = presentation_manager.get_presentation(presentation_id)
            if not (0 <= slide_index < len(pres.slides)):
                return {"error": f"Invalid slide index: {slide_index}"}
            slide = pres.slides[slide_index]
            
            # Use provided bounds or calculate from slide
            if bounds is None:
                bounds = self._get_slide_bounds(presentation_id)
            
            # Calculate cell dimensions
            total_h_gaps = (cols - 1) * gap
            total_v_gaps = (rows - 1) * gap
            cell_width = (bounds.width - total_h_gaps) / cols
            cell_height = (bounds.height - total_v_gaps) / rows
            
            defaults = self._get_default_styles()
            created_shapes = []
            max_elements = rows * cols
            
            # Warn if more elements provided than grid can accommodate
            if len(elements) > max_elements:
                import logging
                logging.getLogger(__name__).warning(
                    f"Grid layout ({rows}x{cols}) can only fit {max_elements} elements, "
                    f"but {len(elements)} were provided. Extra elements will be ignored."
                )
            
            for i, elem_data in enumerate(elements):
                if i >= max_elements:
                    break  # Don't exceed grid capacity
                
                row = i // cols
                col = i % cols
                
                # Calculate position
                left = bounds.left + col * (cell_width + gap)
                top = bounds.top + row * (cell_height + gap)
                
                # Create the element
                elem = self._create_element_from_dict(elem_data, defaults)
                elem.left = left
                elem.top = top
                elem.width = cell_width
                elem.height = cell_height
                
                shape_index = self._add_element_to_slide(slide, elem)
                created_shapes.append({
                    "index": shape_index,
                    "position": {"row": row, "col": col},
                    "bounds": {"left": left, "top": top, "width": cell_width, "height": cell_height}
                })
            
            return {
                "message": f"Created grid layout with {len(created_shapes)} elements ({rows}x{cols})",
                "slide_index": slide_index,
                "layout_type": "grid",
                "shapes": created_shapes
            }
            
        except (ValueError, KeyError) as e:
            return {"error": str(e)}
    
    # =========================================================================
    # LIST LAYOUT
    # =========================================================================
    
    @performance_monitor.track_operation("layout_list")
    def create_list_layout(
        self,
        slide_index: int,
        elements: List[Dict[str, Any]],
        direction: str = "vertical",
        gap: float = 0.2,
        alignment: str = "left",
        bounds: Optional[LayoutBounds] = None,
        presentation_id: Optional[str] = None
    ) -> Dict[str, Any]:
        """
        Create elements arranged in a list (vertical or horizontal).
        
        Args:
            slide_index: Index of the target slide
            elements: List of element dictionaries with 'content' and optional styling
            direction: 'vertical' or 'horizontal'
            gap: Gap between items in inches
            alignment: Alignment of items ('left', 'center', 'right' for vertical;
                      'top', 'middle', 'bottom' for horizontal)
            bounds: Optional custom bounds for the layout area
            presentation_id: Optional presentation ID
            
        Returns:
            Dictionary with confirmation and created shape indices
        """
        try:
            pres = presentation_manager.get_presentation(presentation_id)
            if not (0 <= slide_index < len(pres.slides)):
                return {"error": f"Invalid slide index: {slide_index}"}
            slide = pres.slides[slide_index]
            
            if bounds is None:
                bounds = self._get_slide_bounds(presentation_id)
            
            n_elements = len(elements)
            if n_elements == 0:
                return {"error": "No elements provided"}
            
            defaults = self._get_default_styles()
            created_shapes = []
            
            is_vertical = direction.lower() in ["vertical", "top_to_bottom", "bottom_to_top"]
            
            if is_vertical:
                # Calculate item dimensions for vertical list
                total_gaps = (n_elements - 1) * gap
                item_height = (bounds.height - total_gaps) / n_elements
                item_width = bounds.width * 0.9  # Use 90% width for visual appeal
                
                # Calculate horizontal offset based on alignment
                if alignment.lower() == "center":
                    h_offset = (bounds.width - item_width) / 2
                elif alignment.lower() == "right":
                    h_offset = bounds.width - item_width
                else:  # left
                    h_offset = 0
                
                for i, elem_data in enumerate(elements):
                    elem = self._create_element_from_dict(elem_data, defaults)
                    elem.left = bounds.left + h_offset
                    elem.top = bounds.top + i * (item_height + gap)
                    elem.width = item_width
                    elem.height = item_height
                    
                    shape_index = self._add_element_to_slide(slide, elem)
                    created_shapes.append({
                        "index": shape_index,
                        "position": i,
                        "bounds": {"left": elem.left, "top": elem.top, 
                                  "width": elem.width, "height": elem.height}
                    })
            else:
                # Calculate item dimensions for horizontal list
                total_gaps = (n_elements - 1) * gap
                item_width = (bounds.width - total_gaps) / n_elements
                item_height = bounds.height * 0.6  # Use 60% height for visual appeal
                
                # Calculate vertical offset based on alignment
                if alignment.lower() in ["middle", "center"]:
                    v_offset = (bounds.height - item_height) / 2
                elif alignment.lower() == "bottom":
                    v_offset = bounds.height - item_height
                else:  # top
                    v_offset = 0
                
                for i, elem_data in enumerate(elements):
                    elem = self._create_element_from_dict(elem_data, defaults)
                    elem.left = bounds.left + i * (item_width + gap)
                    elem.top = bounds.top + v_offset
                    elem.width = item_width
                    elem.height = item_height
                    
                    shape_index = self._add_element_to_slide(slide, elem)
                    created_shapes.append({
                        "index": shape_index,
                        "position": i,
                        "bounds": {"left": elem.left, "top": elem.top,
                                  "width": elem.width, "height": elem.height}
                    })
            
            return {
                "message": f"Created {direction} list layout with {len(created_shapes)} elements",
                "slide_index": slide_index,
                "layout_type": "list",
                "direction": direction,
                "shapes": created_shapes
            }
            
        except (ValueError, KeyError) as e:
            return {"error": str(e)}
    
    # =========================================================================
    # HIERARCHY LAYOUT
    # =========================================================================
    
    @performance_monitor.track_operation("layout_hierarchy")
    def create_hierarchy_layout(
        self,
        slide_index: int,
        root: Dict[str, Any],
        level_gap: float = 0.8,
        sibling_gap: float = 0.3,
        show_connectors: bool = True,
        bounds: Optional[LayoutBounds] = None,
        presentation_id: Optional[str] = None
    ) -> Dict[str, Any]:
        """
        Create elements arranged in a hierarchical/tree structure.
        
        Args:
            slide_index: Index of the target slide
            root: Root element with 'content' and optional 'children' list
                  Example: {"content": "CEO", "children": [
                      {"content": "VP Sales", "children": [...]},
                      {"content": "VP Engineering"}
                  ]}
            level_gap: Vertical gap between levels in inches
            sibling_gap: Horizontal gap between siblings in inches
            show_connectors: Whether to draw connecting lines
            bounds: Optional custom bounds for the layout area
            presentation_id: Optional presentation ID
            
        Returns:
            Dictionary with confirmation and created shape indices
        """
        try:
            pres = presentation_manager.get_presentation(presentation_id)
            if not (0 <= slide_index < len(pres.slides)):
                return {"error": f"Invalid slide index: {slide_index}"}
            slide = pres.slides[slide_index]
            
            if bounds is None:
                bounds = self._get_slide_bounds(presentation_id)
            
            defaults = self._get_default_styles()
            
            # Calculate tree structure
            levels = self._flatten_hierarchy(root)
            n_levels = len(levels)
            
            if n_levels == 0:
                return {"error": "Empty hierarchy provided"}
            
            # Calculate dimensions
            max_items_per_level = max(len(level) for level in levels)
            total_v_gaps = (n_levels - 1) * level_gap
            node_height = min(0.8, (bounds.height - total_v_gaps) / n_levels)
            
            created_shapes = []
            node_positions: Dict[int, Tuple[float, float, float, float]] = {}
            
            for level_idx, level_nodes in enumerate(levels):
                n_nodes = len(level_nodes)
                total_h_gaps = (n_nodes - 1) * sibling_gap if n_nodes > 1 else 0
                node_width = min(2.0, (bounds.width - total_h_gaps) / n_nodes)
                
                # Center the level horizontally
                total_level_width = n_nodes * node_width + total_h_gaps
                level_start_x = bounds.left + (bounds.width - total_level_width) / 2
                
                y = bounds.top + level_idx * (node_height + level_gap)
                
                for node_idx, node_data in enumerate(level_nodes):
                    x = level_start_x + node_idx * (node_width + sibling_gap)
                    
                    elem = self._create_element_from_dict(node_data, defaults)
                    elem.left = x
                    elem.top = y
                    elem.width = node_width
                    elem.height = node_height
                    
                    shape_index = self._add_element_to_slide(slide, elem)
                    
                    # Store position for connector drawing (ID assigned by _flatten_hierarchy)
                    node_id = node_data.get("_id")
                    if node_id is not None:
                        node_positions[node_id] = (x, y, node_width, node_height)
                    
                    created_shapes.append({
                        "index": shape_index,
                        "level": level_idx,
                        "position": node_idx,
                        "bounds": {"left": x, "top": y, "width": node_width, "height": node_height}
                    })
            
            # Draw connectors
            connector_shapes = []
            if show_connectors:
                connector_shapes = self._draw_hierarchy_connectors(
                    slide, root, node_positions, defaults
                )
            
            return {
                "message": f"Created hierarchy layout with {len(created_shapes)} nodes and {len(connector_shapes)} connectors",
                "slide_index": slide_index,
                "layout_type": "hierarchy",
                "levels": n_levels,
                "shapes": created_shapes,
                "connectors": connector_shapes
            }
            
        except (ValueError, KeyError) as e:
            return {"error": str(e)}
    
    def _flatten_hierarchy(self, root: Dict[str, Any], level: int = 0,
                          result: Optional[List[List[Dict]]] = None,
                          counter: Optional[List[int]] = None) -> List[List[Dict]]:
        """Flatten hierarchy tree into levels for positioning."""
        if result is None:
            result = []
        if counter is None:
            counter = [0]  # Use list to allow mutation in nested calls
        
        while len(result) <= level:
            result.append([])
        
        # Assign unique sequential ID for reliable node tracking
        node_id = counter[0]
        counter[0] += 1
        root["_id"] = node_id
        result[level].append(root)
        
        children = root.get("children", [])
        for child in children:
            child["_parent_id"] = root["_id"]
            self._flatten_hierarchy(child, level + 1, result, counter)
        
        return result
    
    def _draw_hierarchy_connectors(
        self,
        slide,
        root: Dict[str, Any],
        positions: Dict[int, Tuple[float, float, float, float]],
        defaults: Dict[str, Any]
    ) -> List[Dict[str, Any]]:
        """Draw connecting lines between hierarchy nodes."""
        connectors = []
        
        def draw_connections(node: Dict[str, Any]):
            node_id = node.get("_id")
            if node_id not in positions:
                return
            
            parent_pos = positions[node_id]
            px_center = parent_pos[0] + parent_pos[2] / 2
            py_bottom = parent_pos[1] + parent_pos[3]
            
            children = node.get("children", [])
            for child in children:
                child_id = child.get("_id")
                if child_id in positions:
                    child_pos = positions[child_id]
                    cx_center = child_pos[0] + child_pos[2] / 2
                    cy_top = child_pos[1]
                    
                    # Draw line from parent bottom-center to child top-center
                    ppt_utils.add_line(
                        slide, px_center, py_bottom, cx_center, cy_top,
                        line_color=defaults.get("line_color"),
                        line_width=1.5
                    )
                    connectors.append({
                        "from": {"x": px_center, "y": py_bottom},
                        "to": {"x": cx_center, "y": cy_top}
                    })
                
                draw_connections(child)
        
        draw_connections(root)
        return connectors
    
    # =========================================================================
    # FLOW LAYOUT
    # =========================================================================
    
    @performance_monitor.track_operation("layout_flow")
    def create_flow_layout(
        self,
        slide_index: int,
        steps: List[Dict[str, Any]],
        direction: str = "horizontal",
        gap: float = 0.4,
        show_connectors: bool = True,
        connector_style: str = "arrow",
        bounds: Optional[LayoutBounds] = None,
        presentation_id: Optional[str] = None
    ) -> Dict[str, Any]:
        """
        Create elements arranged as a flow/process diagram.
        
        Args:
            slide_index: Index of the target slide
            steps: List of step dictionaries with 'content' and optional styling
            direction: 'horizontal' (left to right) or 'vertical' (top to bottom)
            gap: Gap between steps in inches (includes space for connectors)
            show_connectors: Whether to draw connecting arrows/lines
            connector_style: 'arrow', 'line', or 'none'
            bounds: Optional custom bounds for the layout area
            presentation_id: Optional presentation ID
            
        Returns:
            Dictionary with confirmation and created shape indices
        """
        try:
            pres = presentation_manager.get_presentation(presentation_id)
            if not (0 <= slide_index < len(pres.slides)):
                return {"error": f"Invalid slide index: {slide_index}"}
            slide = pres.slides[slide_index]
            
            if bounds is None:
                bounds = self._get_slide_bounds(presentation_id)
            
            n_steps = len(steps)
            if n_steps == 0:
                return {"error": "No steps provided"}
            
            defaults = self._get_default_styles()
            created_shapes = []
            step_positions = []
            
            is_horizontal = direction.lower() in ["horizontal", "left_to_right"]
            
            if is_horizontal:
                # Calculate step dimensions for horizontal flow
                total_gaps = (n_steps - 1) * gap
                step_width = min(2.0, (bounds.width - total_gaps) / n_steps)
                step_height = min(1.5, bounds.height * 0.5)
                
                # Center vertically
                y = bounds.top + (bounds.height - step_height) / 2
                
                for i, step_data in enumerate(steps):
                    x = bounds.left + i * (step_width + gap)
                    
                    elem = self._create_element_from_dict(step_data, defaults)
                    # Use rounded rectangle for flow steps by default
                    if elem.element_type == "textbox":
                        elem.element_type = "shape"
                        elem.shape_type = step_data.get("shape_type", "rounded_rectangle")
                    
                    elem.left = x
                    elem.top = y
                    elem.width = step_width
                    elem.height = step_height
                    
                    shape_index = self._add_element_to_slide(slide, elem)
                    step_positions.append((x, y, step_width, step_height))
                    
                    created_shapes.append({
                        "index": shape_index,
                        "step": i,
                        "bounds": {"left": x, "top": y, "width": step_width, "height": step_height}
                    })
            else:
                # Calculate step dimensions for vertical flow
                total_gaps = (n_steps - 1) * gap
                step_height = min(1.0, (bounds.height - total_gaps) / n_steps)
                step_width = min(3.0, bounds.width * 0.6)
                
                # Center horizontally
                x = bounds.left + (bounds.width - step_width) / 2
                
                for i, step_data in enumerate(steps):
                    y = bounds.top + i * (step_height + gap)
                    
                    elem = self._create_element_from_dict(step_data, defaults)
                    if elem.element_type == "textbox":
                        elem.element_type = "shape"
                        elem.shape_type = step_data.get("shape_type", "rounded_rectangle")
                    
                    elem.left = x
                    elem.top = y
                    elem.width = step_width
                    elem.height = step_height
                    
                    shape_index = self._add_element_to_slide(slide, elem)
                    step_positions.append((x, y, step_width, step_height))
                    
                    created_shapes.append({
                        "index": shape_index,
                        "step": i,
                        "bounds": {"left": x, "top": y, "width": step_width, "height": step_height}
                    })
            
            # Draw connectors between steps
            connector_shapes = []
            if show_connectors and connector_style != "none" and n_steps > 1:
                connector_shapes = self._draw_flow_connectors(
                    slide, step_positions, direction, connector_style, defaults
                )
            
            return {
                "message": f"Created {direction} flow layout with {len(created_shapes)} steps",
                "slide_index": slide_index,
                "layout_type": "flow",
                "direction": direction,
                "shapes": created_shapes,
                "connectors": connector_shapes
            }
            
        except (ValueError, KeyError) as e:
            return {"error": str(e)}
    
    def _draw_flow_connectors(
        self,
        slide,
        positions: List[Tuple[float, float, float, float]],
        direction: str,
        style: str,
        defaults: Dict[str, Any]
    ) -> List[Dict[str, Any]]:
        """Draw connecting lines/arrows between flow steps."""
        from pptx.util import Inches
        from pptx.enum.shapes import MSO_SHAPE
        
        connectors = []
        is_horizontal = direction.lower() in ["horizontal", "left_to_right"]
        
        for i in range(len(positions) - 1):
            pos1 = positions[i]
            pos2 = positions[i + 1]
            
            if is_horizontal:
                # Draw from right edge of step1 to left edge of step2
                x1 = pos1[0] + pos1[2]  # Right edge
                y1 = pos1[1] + pos1[3] / 2  # Vertical center
                x2 = pos2[0]  # Left edge
                y2 = pos2[1] + pos2[3] / 2  # Vertical center
            else:
                # Draw from bottom edge of step1 to top edge of step2
                x1 = pos1[0] + pos1[2] / 2  # Horizontal center
                y1 = pos1[1] + pos1[3]  # Bottom edge
                x2 = pos2[0] + pos2[2] / 2  # Horizontal center
                y2 = pos2[1]  # Top edge
            
            if style == "arrow":
                # Use proper arrow shapes instead of triangles
                # Calculate arrow dimensions and position
                if is_horizontal:
                    arrow_width = x2 - x1
                    arrow_height = 0.25
                    arrow_left = x1
                    arrow_top = y1 - arrow_height / 2
                    # Use RIGHT_ARROW shape for horizontal flow
                    arrow_shape = slide.shapes.add_shape(
                        MSO_SHAPE.RIGHT_ARROW,
                        Inches(arrow_left), Inches(arrow_top),
                        Inches(arrow_width), Inches(arrow_height)
                    )
                else:
                    arrow_width = 0.25
                    arrow_height = y2 - y1
                    arrow_left = x1 - arrow_width / 2
                    arrow_top = y1
                    # Use DOWN_ARROW shape for vertical flow
                    arrow_shape = slide.shapes.add_shape(
                        MSO_SHAPE.DOWN_ARROW,
                        Inches(arrow_left), Inches(arrow_top),
                        Inches(arrow_width), Inches(arrow_height)
                    )
                
                # Style the arrow
                line_color = defaults.get("line_color", [0, 0, 0])
                if arrow_shape.fill:
                    arrow_shape.fill.solid()
                    from pptx.dml.color import RGBColor
                    arrow_shape.fill.fore_color.rgb = RGBColor(*line_color)
                if arrow_shape.line:
                    arrow_shape.line.fill.background()  # No outline
            else:
                # Just draw a line
                ppt_utils.add_line(
                    slide, x1, y1, x2, y2,
                    line_color=defaults.get("line_color"),
                    line_width=1.5
                )
            
            connectors.append({
                "from": {"x": x1, "y": y1},
                "to": {"x": x2, "y": y2},
                "style": style
            })
        
        return connectors
    
    # =========================================================================
    # HELPER METHODS
    # =========================================================================
    
    def _create_element_from_dict(
        self,
        data: Dict[str, Any],
        defaults: Dict[str, Any]
    ) -> LayoutElement:
        """Create a LayoutElement from a dictionary with defaults."""
        return LayoutElement(
            content=str(data.get("content", "")),
            element_type=data.get("element_type", "textbox"),
            shape_type=data.get("shape_type"),
            fill_color=data.get("fill_color", defaults.get("fill_color")),
            line_color=data.get("line_color", defaults.get("line_color")),
            font_size=data.get("font_size", defaults.get("font_size")),
            font_name=data.get("font_name", defaults.get("font_name")),
            bold=data.get("bold"),
            text_color=data.get("text_color", defaults.get("text_color")),
            alignment=data.get("alignment", "center")
        )
    
    def _add_element_to_slide(self, slide, elem: LayoutElement) -> int:
        """Add a layout element to a slide and return the shape index."""
        import logging
        logger = logging.getLogger(__name__)
        
        if elem.element_type == "shape" and elem.shape_type:
            # Add shape with text
            shape = ppt_utils.add_shape(
                slide, elem.shape_type,
                elem.left, elem.top, elem.width, elem.height,
                fill_color=elem.fill_color,
                line_color=elem.line_color
            )
            # Add text to the shape
            if elem.content:
                if hasattr(shape, 'text_frame'):
                    tf = shape.text_frame
                    tf.word_wrap = True
                    p = tf.paragraphs[0]
                    p.text = elem.content
                    if elem.alignment:
                        from pptx.enum.text import PP_ALIGN
                        # Validate alignment value
                        alignment_upper = elem.alignment.upper()
                        if hasattr(PP_ALIGN, alignment_upper):
                            p.alignment = getattr(PP_ALIGN, alignment_upper)
                        else:
                            logger.warning(f"Invalid alignment '{elem.alignment}', using CENTER")
                            p.alignment = PP_ALIGN.CENTER
                    font = p.font
                    if elem.font_size:
                        from pptx.util import Pt
                        font.size = Pt(elem.font_size)
                    if elem.font_name:
                        font.name = elem.font_name
                    if elem.bold is not None:
                        font.bold = elem.bold
                    if elem.text_color:
                        from pptx.dml.color import RGBColor
                        font.color.rgb = RGBColor(*elem.text_color)
                else:
                    logger.warning(f"Shape type '{elem.shape_type}' does not support text_frame, text content ignored")
        else:
            # Add textbox
            ppt_utils.add_textbox(
                slide, elem.left, elem.top, elem.width, elem.height,
                elem.content,
                font_size=elem.font_size,
                font_name=elem.font_name,
                bold=elem.bold,
                color=elem.text_color,
                alignment=elem.alignment
            )
        
        return len(slide.shapes) - 1


# Global instance
layout_manager = LayoutEngine()
