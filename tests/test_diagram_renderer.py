#!/usr/bin/env python3
"""
Tests for the diagram renderer module.

Tests the rendering of parsed diagrams (from Mermaid/PlantUML) as editable
PowerPoint vector shapes.
"""

import sys
import os
sys.path.insert(0, os.path.dirname(os.path.dirname(os.path.abspath(__file__))))

import unittest
from unittest.mock import patch, MagicMock
from pptx import Presentation

from diagram_renderer import (
    DiagramRenderer, DiagramStyle, get_default_diagram_style,
    diagram_renderer
)
from diagram_parser import (
    ParsedDiagram, DiagramNode, DiagramEdge,
    NodeShape, Direction, EdgeType, EdgeStyle
)
from layout_manager import LayoutBounds


# Slide layout index for blank slides in default PowerPoint templates
BLANK_SLIDE_LAYOUT_INDEX = 6


class TestDiagramRendererInitialization(unittest.TestCase):
    """Tests for DiagramRenderer initialization."""
    
    def test_initialization(self):
        """Test that DiagramRenderer initializes correctly."""
        renderer = DiagramRenderer()
        self.assertIsNotNone(renderer._parser)
    
    def test_global_instance_exists(self):
        """Test that global diagram_renderer instance exists."""
        self.assertIsInstance(diagram_renderer, DiagramRenderer)


class TestDiagramStyle(unittest.TestCase):
    """Tests for DiagramStyle dataclass."""
    
    def test_default_values(self):
        """Test default values for DiagramStyle."""
        style = DiagramStyle()
        
        self.assertIsNone(style.default_fill_color)
        self.assertIsNone(style.default_text_color)
        self.assertIsNone(style.default_line_color)
        self.assertEqual(style.connector_width, 1.5)
        self.assertEqual(style.font_size, 14)
        self.assertFalse(style.bold)
    
    def test_custom_values(self):
        """Test creating DiagramStyle with custom values."""
        style = DiagramStyle(
            default_fill_color=[255, 0, 0],
            default_text_color=[255, 255, 255],
            connector_width=2.0,
            font_size=18,
            bold=True
        )
        
        self.assertEqual(style.default_fill_color, [255, 0, 0])
        self.assertEqual(style.default_text_color, [255, 255, 255])
        self.assertEqual(style.connector_width, 2.0)
        self.assertEqual(style.font_size, 18)
        self.assertTrue(style.bold)


class TestGetDefaultDiagramStyle(unittest.TestCase):
    """Tests for get_default_diagram_style function."""
    
    @patch('diagram_renderer.template_manager')
    def test_get_default_diagram_style(self, mock_tm):
        """Test getting default diagram style from template."""
        mock_tm.get_default_font_settings.return_value = {
            "body_font_name": "Arial",
            "body_font_size": 16
        }
        mock_tm.get_default_color_settings.return_value = {
            "accent_1": (100, 150, 200),
            "text_1": (50, 50, 50)
        }
        
        style = get_default_diagram_style()
        
        self.assertIsInstance(style, DiagramStyle)
        self.assertEqual(style.font_name, "Arial")
        self.assertEqual(style.default_fill_color, [100, 150, 200])


class TestRenderMermaid(unittest.TestCase):
    """Tests for render_mermaid method."""
    
    def setUp(self):
        """Set up test fixtures."""
        self.renderer = DiagramRenderer()
        self.pres = Presentation()
        self.pres.slides.add_slide(self.pres.slide_layouts[BLANK_SLIDE_LAYOUT_INDEX])
    
    @patch('diagram_renderer.layout_manager')
    @patch('diagram_renderer.presentation_manager')
    def test_render_mermaid_simple_flow(self, mock_pm, mock_lm):
        """Test rendering a simple Mermaid flow diagram."""
        mock_pm.get_presentation.return_value = self.pres
        mock_lm.create_flow_layout.return_value = {
            "message": "Created flow layout",
            "shapes": []
        }
        mock_lm._get_slide_bounds.return_value = LayoutBounds()
        
        mermaid_code = """
graph TD
    A[Start] --> B[End]
"""
        
        result = self.renderer.render_mermaid(
            slide_index=0,
            mermaid_code=mermaid_code
        )
        
        self.assertNotIn("error", result)
    
    @patch('diagram_renderer.presentation_manager')
    def test_render_mermaid_invalid_slide(self, mock_pm):
        """Test rendering with invalid slide index."""
        mock_pm.get_presentation.return_value = self.pres
        
        result = self.renderer.render_mermaid(
            slide_index=99,
            mermaid_code="graph TD\n    A --> B"
        )
        
        self.assertIn("error", result)
    
    @patch('diagram_renderer.presentation_manager')
    def test_render_mermaid_parse_error(self, mock_pm):
        """Test rendering with invalid Mermaid code."""
        mock_pm.get_presentation.return_value = self.pres
        
        result = self.renderer.render_mermaid(
            slide_index=0,
            mermaid_code=""  # Empty code should fail
        )
        
        self.assertIn("error", result)
    
    @patch('diagram_renderer.layout_manager')
    @patch('diagram_renderer.presentation_manager')
    def test_render_mermaid_with_custom_style(self, mock_pm, mock_lm):
        """Test rendering with custom style."""
        mock_pm.get_presentation.return_value = self.pres
        mock_lm.create_flow_layout.return_value = {
            "message": "Created flow layout",
            "shapes": []
        }
        mock_lm._get_slide_bounds.return_value = LayoutBounds()
        
        custom_style = DiagramStyle(
            default_fill_color=[255, 0, 0],
            font_size=20
        )
        
        result = self.renderer.render_mermaid(
            slide_index=0,
            mermaid_code="graph TD\n    A --> B",
            style=custom_style
        )
        
        self.assertNotIn("error", result)


class TestRenderPlantUML(unittest.TestCase):
    """Tests for render_plantuml method."""
    
    def setUp(self):
        """Set up test fixtures."""
        self.renderer = DiagramRenderer()
        self.pres = Presentation()
        self.pres.slides.add_slide(self.pres.slide_layouts[BLANK_SLIDE_LAYOUT_INDEX])
    
    @patch('diagram_renderer.layout_manager')
    @patch('diagram_renderer.presentation_manager')
    def test_render_plantuml_simple(self, mock_pm, mock_lm):
        """Test rendering a simple PlantUML diagram."""
        mock_pm.get_presentation.return_value = self.pres
        mock_lm.create_flow_layout.return_value = {
            "message": "Created flow layout",
            "shapes": []
        }
        mock_lm._get_slide_bounds.return_value = LayoutBounds()
        
        plantuml_code = """
@startuml
start
:Action;
stop
@enduml
"""
        
        result = self.renderer.render_plantuml(
            slide_index=0,
            plantuml_code=plantuml_code
        )
        
        self.assertNotIn("error", result)
    
    @patch('diagram_renderer.presentation_manager')
    def test_render_plantuml_invalid_slide(self, mock_pm):
        """Test rendering with invalid slide index."""
        mock_pm.get_presentation.return_value = self.pres
        
        result = self.renderer.render_plantuml(
            slide_index=99,
            plantuml_code="start\nstop"
        )
        
        self.assertIn("error", result)


class TestRenderAuto(unittest.TestCase):
    """Tests for render_auto method."""
    
    def setUp(self):
        """Set up test fixtures."""
        self.renderer = DiagramRenderer()
        self.pres = Presentation()
        self.pres.slides.add_slide(self.pres.slide_layouts[BLANK_SLIDE_LAYOUT_INDEX])
    
    @patch('diagram_renderer.layout_manager')
    @patch('diagram_renderer.presentation_manager')
    def test_render_auto_mermaid(self, mock_pm, mock_lm):
        """Test auto-detecting and rendering Mermaid code."""
        mock_pm.get_presentation.return_value = self.pres
        mock_lm.create_flow_layout.return_value = {
            "message": "Created flow layout",
            "shapes": []
        }
        mock_lm._get_slide_bounds.return_value = LayoutBounds()
        
        result = self.renderer.render_auto(
            slide_index=0,
            diagram_code="graph TD\n    A --> B"
        )
        
        self.assertNotIn("error", result)
        self.assertEqual(result.get("detected_type"), "mermaid")
    
    @patch('diagram_renderer.layout_manager')
    @patch('diagram_renderer.presentation_manager')
    def test_render_auto_plantuml(self, mock_pm, mock_lm):
        """Test auto-detecting and rendering PlantUML code."""
        mock_pm.get_presentation.return_value = self.pres
        mock_lm.create_flow_layout.return_value = {
            "message": "Created flow layout",
            "shapes": []
        }
        mock_lm._get_slide_bounds.return_value = LayoutBounds()
        
        result = self.renderer.render_auto(
            slide_index=0,
            diagram_code="@startuml\nstart\nstop\n@enduml"
        )
        
        self.assertNotIn("error", result)
        self.assertEqual(result.get("detected_type"), "plantuml")


class TestDiagramTypeDetection(unittest.TestCase):
    """Tests for diagram type detection helper methods."""
    
    def setUp(self):
        """Set up test fixtures."""
        self.renderer = DiagramRenderer()
    
    def test_is_linear_flow_true(self):
        """Test detection of linear flow diagram."""
        from diagram_parser import DiagramType
        diagram = ParsedDiagram(
            diagram_type=DiagramType.FLOWCHART,
            direction=Direction.LEFT_RIGHT,
            nodes=[
                DiagramNode(id="A", label="A"),
                DiagramNode(id="B", label="B"),
                DiagramNode(id="C", label="C")
            ],
            edges=[
                DiagramEdge(source="A", target="B"),
                DiagramEdge(source="B", target="C")
            ]
        )
        
        result = self.renderer._is_linear_flow(diagram)
        
        self.assertTrue(result)
    
    def test_is_linear_flow_false_branching(self):
        """Test detection of non-linear (branching) diagram."""
        from diagram_parser import DiagramType
        diagram = ParsedDiagram(
            diagram_type=DiagramType.FLOWCHART,
            direction=Direction.TOP_DOWN,
            nodes=[
                DiagramNode(id="A", label="A"),
                DiagramNode(id="B", label="B"),
                DiagramNode(id="C", label="C")
            ],
            edges=[
                DiagramEdge(source="A", target="B"),
                DiagramEdge(source="A", target="C")  # Branching
            ]
        )
        
        result = self.renderer._is_linear_flow(diagram)
        
        self.assertFalse(result)
    
    def test_is_linear_flow_no_edges(self):
        """Test linear flow with no edges returns True."""
        from diagram_parser import DiagramType
        diagram = ParsedDiagram(
            diagram_type=DiagramType.FLOWCHART,
            direction=Direction.TOP_DOWN,
            nodes=[DiagramNode(id="A", label="A")],
            edges=[]
        )
        
        result = self.renderer._is_linear_flow(diagram)
        
        self.assertTrue(result)
    
    def test_is_hierarchical_true(self):
        """Test detection of hierarchical diagram."""
        from diagram_parser import DiagramType
        diagram = ParsedDiagram(
            diagram_type=DiagramType.HIERARCHY,
            direction=Direction.TOP_DOWN,
            nodes=[
                DiagramNode(id="root", label="Root"),
                DiagramNode(id="child1", label="Child 1"),
                DiagramNode(id="child2", label="Child 2")
            ],
            edges=[
                DiagramEdge(source="root", target="child1"),
                DiagramEdge(source="root", target="child2")
            ]
        )
        
        result = self.renderer._is_hierarchical(diagram)
        
        self.assertTrue(result)
    
    def test_is_hierarchical_false_no_edges(self):
        """Test hierarchical detection with no edges returns False."""
        from diagram_parser import DiagramType
        diagram = ParsedDiagram(
            diagram_type=DiagramType.FLOWCHART,
            direction=Direction.TOP_DOWN,
            nodes=[DiagramNode(id="A", label="A")],
            edges=[]
        )
        
        result = self.renderer._is_hierarchical(diagram)
        
        self.assertFalse(result)
    
    def test_is_hierarchical_false_multiple_roots(self):
        """Test hierarchical detection with multiple roots."""
        from diagram_parser import DiagramType
        diagram = ParsedDiagram(
            diagram_type=DiagramType.FLOWCHART,
            direction=Direction.TOP_DOWN,
            nodes=[
                DiagramNode(id="A", label="A"),
                DiagramNode(id="B", label="B"),
                DiagramNode(id="C", label="C")
            ],
            edges=[
                DiagramEdge(source="A", target="C"),
                DiagramEdge(source="B", target="C")  # Two roots
            ]
        )
        
        result = self.renderer._is_hierarchical(diagram)
        
        self.assertFalse(result)


class TestNodeOrdering(unittest.TestCase):
    """Tests for node ordering in flow diagrams."""
    
    def setUp(self):
        """Set up test fixtures."""
        self.renderer = DiagramRenderer()
    
    def test_order_nodes_for_flow(self):
        """Test ordering nodes for flow diagram."""
        from diagram_parser import DiagramType
        diagram = ParsedDiagram(
            diagram_type=DiagramType.FLOWCHART,
            direction=Direction.LEFT_RIGHT,
            nodes=[
                DiagramNode(id="C", label="C"),
                DiagramNode(id="A", label="A"),
                DiagramNode(id="B", label="B")
            ],
            edges=[
                DiagramEdge(source="A", target="B"),
                DiagramEdge(source="B", target="C")
            ]
        )
        
        ordered = self.renderer._order_nodes_for_flow(diagram)
        
        # Should be ordered A -> B -> C
        self.assertEqual(len(ordered), 3)
        self.assertEqual(ordered[0].id, "A")
        self.assertEqual(ordered[1].id, "B")
        self.assertEqual(ordered[2].id, "C")
    
    def test_order_nodes_no_edges(self):
        """Test ordering nodes when no edges exist."""
        from diagram_parser import DiagramType
        diagram = ParsedDiagram(
            diagram_type=DiagramType.FLOWCHART,
            direction=Direction.TOP_DOWN,
            nodes=[
                DiagramNode(id="A", label="A"),
                DiagramNode(id="B", label="B")
            ],
            edges=[]
        )
        
        ordered = self.renderer._order_nodes_for_flow(diagram)
        
        # Should return all nodes
        self.assertEqual(len(ordered), 2)
    
    def test_order_nodes_disconnected(self):
        """Test ordering with disconnected nodes."""
        from diagram_parser import DiagramType
        diagram = ParsedDiagram(
            diagram_type=DiagramType.FLOWCHART,
            direction=Direction.LEFT_RIGHT,
            nodes=[
                DiagramNode(id="A", label="A"),
                DiagramNode(id="B", label="B"),
                DiagramNode(id="X", label="X")  # Disconnected
            ],
            edges=[
                DiagramEdge(source="A", target="B")
            ]
        )
        
        ordered = self.renderer._order_nodes_for_flow(diagram)
        
        # Should include all nodes
        self.assertEqual(len(ordered), 3)
        node_ids = [n.id for n in ordered]
        self.assertIn("X", node_ids)


class TestBuildHierarchyTree(unittest.TestCase):
    """Tests for _build_hierarchy_tree method."""
    
    def setUp(self):
        """Set up test fixtures."""
        self.renderer = DiagramRenderer()
        self.style = DiagramStyle(
            default_fill_color=[100, 100, 100],
            default_text_color=[255, 255, 255]
        )
    
    def test_build_hierarchy_tree_simple(self):
        """Test building simple hierarchy tree."""
        from diagram_parser import DiagramType
        diagram = ParsedDiagram(
            diagram_type=DiagramType.HIERARCHY,
            direction=Direction.TOP_DOWN,
            nodes=[
                DiagramNode(id="root", label="Root"),
                DiagramNode(id="child", label="Child")
            ],
            edges=[
                DiagramEdge(source="root", target="child")
            ]
        )
        
        tree = self.renderer._build_hierarchy_tree(diagram, self.style)
        
        self.assertIsNotNone(tree)
        self.assertEqual(tree["content"], "Root")
        self.assertIn("children", tree)
        self.assertEqual(len(tree["children"]), 1)
    
    def test_build_hierarchy_tree_no_nodes(self):
        """Test building hierarchy tree with no nodes."""
        from diagram_parser import DiagramType
        diagram = ParsedDiagram(
            diagram_type=DiagramType.FLOWCHART,
            direction=Direction.TOP_DOWN,
            nodes=[],
            edges=[]
        )
        
        tree = self.renderer._build_hierarchy_tree(diagram, self.style)
        
        self.assertIsNone(tree)
    
    def test_build_hierarchy_tree_no_clear_root(self):
        """Test building hierarchy tree with no clear root."""
        from diagram_parser import DiagramType
        diagram = ParsedDiagram(
            diagram_type=DiagramType.FLOWCHART,
            direction=Direction.TOP_DOWN,
            nodes=[
                DiagramNode(id="A", label="A"),
                DiagramNode(id="B", label="B")
            ],
            edges=[
                DiagramEdge(source="A", target="B"),
                DiagramEdge(source="B", target="A")  # Cycle, no root
            ]
        )
        
        tree = self.renderer._build_hierarchy_tree(diagram, self.style)
        
        # Should return None when no clear root
        self.assertIsNone(tree)


class TestNodeToElement(unittest.TestCase):
    """Tests for _node_to_element method."""
    
    def setUp(self):
        """Set up test fixtures."""
        self.renderer = DiagramRenderer()
        self.style = DiagramStyle(
            default_fill_color=[100, 100, 100],
            default_text_color=[255, 255, 255],
            default_line_color=[0, 0, 0],
            font_name="Arial",
            font_size=14,
            bold=True
        )
    
    def test_node_to_element_basic(self):
        """Test converting basic node to element."""
        node = DiagramNode(
            id="test",
            label="Test Node",
            shape=NodeShape.RECTANGLE
        )
        
        element = self.renderer._node_to_element(node, self.style)
        
        self.assertEqual(element["content"], "Test Node")
        self.assertEqual(element["element_type"], "shape")
        self.assertEqual(element["shape_type"], "rectangle")
        self.assertEqual(element["font_size"], 14)
        self.assertTrue(element["bold"])
    
    def test_node_to_element_rounded_rectangle(self):
        """Test converting rounded rectangle node."""
        node = DiagramNode(
            id="test",
            label="Rounded",
            shape=NodeShape.ROUNDED_RECTANGLE
        )
        
        element = self.renderer._node_to_element(node, self.style)
        
        self.assertEqual(element["shape_type"], "rounded_rectangle")
    
    def test_node_to_element_diamond(self):
        """Test converting diamond node."""
        node = DiagramNode(
            id="test",
            label="Decision",
            shape=NodeShape.DIAMOND
        )
        
        element = self.renderer._node_to_element(node, self.style)
        
        self.assertEqual(element["shape_type"], "diamond")
    
    def test_node_to_element_circle(self):
        """Test converting circle node."""
        node = DiagramNode(
            id="test",
            label="Circle",
            shape=NodeShape.CIRCLE
        )
        
        element = self.renderer._node_to_element(node, self.style)
        
        self.assertEqual(element["shape_type"], "oval")
    
    def test_node_to_element_with_node_colors(self):
        """Test node-specific colors override style defaults."""
        node = DiagramNode(
            id="test",
            label="Custom Colors",
            shape=NodeShape.RECTANGLE,
            fill_color=[255, 0, 0],
            text_color=[0, 255, 0]
        )
        
        element = self.renderer._node_to_element(node, self.style)
        
        self.assertEqual(element["fill_color"], [255, 0, 0])
        self.assertEqual(element["text_color"], [0, 255, 0])
    
    def test_node_to_element_uses_style_defaults(self):
        """Test style defaults are used when node has no colors."""
        node = DiagramNode(
            id="test",
            label="Default Colors",
            shape=NodeShape.RECTANGLE
        )
        
        element = self.renderer._node_to_element(node, self.style)
        
        self.assertEqual(element["fill_color"], [100, 100, 100])
        self.assertEqual(element["text_color"], [255, 255, 255])


class TestShapeMapping(unittest.TestCase):
    """Tests for node shape to PowerPoint shape type mapping."""
    
    def setUp(self):
        """Set up test fixtures."""
        self.renderer = DiagramRenderer()
        self.style = DiagramStyle()
    
    def test_all_shapes_map_correctly(self):
        """Test that all NodeShape values map to valid shape types."""
        shape_tests = [
            (NodeShape.RECTANGLE, "rectangle"),
            (NodeShape.ROUNDED_RECTANGLE, "rounded_rectangle"),
            (NodeShape.DIAMOND, "diamond"),
            (NodeShape.CIRCLE, "oval"),
            (NodeShape.STADIUM, "rounded_rectangle"),
            (NodeShape.HEXAGON, "hexagon"),
            (NodeShape.DATABASE, "flowchart_document"),
        ]
        
        for node_shape, expected_type in shape_tests:
            node = DiagramNode(
                id="test",
                label="Test",
                shape=node_shape
            )
            element = self.renderer._node_to_element(node, self.style)
            self.assertEqual(
                element["shape_type"],
                expected_type,
                f"Failed for {node_shape}"
            )


if __name__ == '__main__':
    unittest.main()
